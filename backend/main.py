
import io
import json 
import PyPDF2
import pytesseract
from PIL import Image
try:
    from pdf2image import convert_from_path
except Exception:
    convert_from_path = None
# Optional PDF renderer that doesn't require Poppler (great for Windows)
import pypdfium2 as pdfium
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None
try:
    import pdfplumber
except Exception:
    pdfplumber = None
# Existing pure-Python PDF text extractor (no extra system deps)
try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None
import docx
import openpyxl
import pptx
import tempfile
import shutil
import glob
try:
    import winreg  # Windows registry, used to detect Tesseract install dir
except Exception:
    winreg = None
from fastapi import FastAPI, Request, UploadFile, File, Form, WebSocket, WebSocketDisconnect
from fastapi.responses import JSONResponse
from llm_inference import generate_response, summarize_text, answer_question, analyze_pdf, get_settings, set_settings
from llm_inference import get_model_name as llm_get_model_name, set_model_name as llm_set_model_name
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List
import re
import os
import sqlite3
from sentence_transformers import SentenceTransformer
from lingua import Language, LanguageDetectorBuilder
from fastapi.staticfiles import StaticFiles
import torch
import logging
import re, threading, time, unicodedata, asyncio, uuid
import soundfile as sf
import numpy as np
try:
    from TTS.api import TTS  # Provided by the coqui-tts package
except Exception:
    TTS = None  # type: ignore
# Configure root logger so your log messages show up in uvicorn console
logging.basicConfig(
    level=logging.DEBUG,   # or INFO
    format="%(asctime)s - %(levelname)s - %(message)s"
)
# Suppress noisy UIAutomation element release spam that floods console
class _SuppressUIAReleaseFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:  # type: ignore[override]
        try:
            msg = record.getMessage()
        except Exception:
            return True
        if msg.startswith('Release <POINTER(IUIAutomationElement)'):
            return False
        return True

logging.getLogger().addFilter(_SuppressUIAReleaseFilter())
# Quiet specific noisy libraries if present
for _noisy in ('comtypes', 'UIAutomationClient', 'uiautomation'):
    logging.getLogger(_noisy).setLevel(logging.WARNING)

app = FastAPI()

# --- Live LAN sharing state ---
_LIVE_LAN_ENABLED = True  # default on so existing functionality works
_LIVE_LAN_CLIENTS: dict[str, float] = {}  # ip -> last_seen timestamp
_LIVE_LAN_CLIENT_TTL = 300  # 5 minutes

# Allow frontend to access backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --------- General Health Check ---------
@app.get('/api/health')
def general_health():
    """Summarize backend health, dependency & feature readiness matrix for UI.

    Adds a structured "features" list where each feature declares its required
    components (dependencies) and whether it's ready to use. This allows the
    frontend to quickly show what parts of the product are operational and
    what is missing / optional.
    """
    import sys, hashlib, shutil, importlib, sqlite3 as _sqlite3, time as _t

    def _dep(id_: str, ok: bool, detail: str = '', optional: bool = False, group: str | None = None):
        return {
            'id': id_,
            'ok': bool(ok),
            'detail': detail or None,
            'optional': optional,
            'group': group,
        }

    # Base skeleton
    info: dict[str, any] = {
        'ok': True,
        'python': sys.version,
        'os': os.name,
        'platform': sys.platform,
        'timestamp': int(_t.time()),
        'packages': {},
        'ocr': {
            'tesseract_cmd': TESSERACT_CMD,
            'poppler_path': POPPLER_PATH,
        },
        'automation': {},
        'features': [],
    }

    # --- Low level dependency introspection ---
    # torch
    try:
        import torch as _torch
        info['packages']['torch'] = {
            'version': getattr(_torch, '__version__', None),
            'cuda_available': bool(getattr(_torch, 'cuda', None) and _torch.cuda.is_available()),
        }
    except Exception as e:
        info['packages']['torch'] = {'error': str(e)}
    # sentence-transformers
    try:
        import sentence_transformers as _st
        info['packages']['sentence_transformers'] = {'version': getattr(_st, '__version__', None)}
    except Exception as e:
        info['packages']['sentence_transformers'] = {'error': str(e)}
    # TTS
    try:
        import TTS as _coqui
        info['packages']['TTS'] = {'version': getattr(_coqui, '__version__', None)}
    except Exception as e:
        info['packages']['TTS'] = {'error': str(e)}
    # pytesseract version best-effort
    try:
        v = str(pytesseract.get_tesseract_version()) if pytesseract is not None else None
    except Exception:
        v = None
    info['ocr']['tesseract_version'] = v

    # Word COM availability
    try:
        import win32com.client  # type: ignore
        info['automation']['win32com'] = True
    except Exception as e:
        info['automation']['win32com'] = False
        info['automation']['win32com_error'] = str(e)
    # comtypes
    try:
        import comtypes  # type: ignore
        info['automation']['comtypes'] = True
    except Exception as e:
        info['automation']['comtypes'] = False
        info['automation']['comtypes_error'] = str(e)

    # Basic DB presence
    db_exists = False
    db_path = os.path.join(os.path.dirname(__file__), 'chat_embeddings.db')
    try:
        db_exists = os.path.isfile(db_path)
        if db_exists:
            info['db'] = {'path': db_path, 'size_bytes': os.path.getsize(db_path)}
        else:
            info['db'] = {'path': db_path, 'missing': True}
    except Exception as e:
        info['db'] = {'error': str(e)}

    # Try to read minimal metadata from DB (optional)
    db_tables = []
    if db_exists:
        try:
            conn = _sqlite3.connect(db_path)
            c = conn.cursor()
            c.execute("SELECT name FROM sqlite_master WHERE type='table'")
            db_tables = [r[0] for r in c.fetchall() if r and r[0]]
            conn.close()
        except Exception:
            db_tables = []
    info['db']['tables'] = db_tables

    # VS Code resolution attempt (lightweight)
    vscode_candidates = []
    resolved_vscode = None
    try:
        if os.name == 'nt':
            # Check env override
            env_bin = os.environ.get('POWER_VSCODE_BIN')
            if env_bin and os.path.isfile(env_bin):
                resolved_vscode = env_bin
            else:
                which_code = shutil.which('code')
                if which_code:
                    resolved_vscode = which_code
        else:
            which_code = shutil.which('code')
            if which_code:
                resolved_vscode = which_code
    except Exception:
        pass
    if resolved_vscode:
        vscode_candidates.append(resolved_vscode)
    info['automation']['vscode_candidate'] = resolved_vscode

    # CUA adapter availability
    cua_status = {}
    try:
        try:
            import backend.cua_adapter as _cua  # type: ignore
        except Exception:
            import cua_adapter as _cua  # type: ignore
        if hasattr(_cua, 'cua_runtime_status'):
            cua_status = _cua.cua_runtime_status()
    except Exception as e:
        cua_status = {'error': str(e)}
    info['automation']['cua'] = cua_status

    # Robust Word COM probe encapsulated to avoid CoInitialize errors & variable scoping issues
    def _probe_word_com() -> tuple[bool, str]:
        if os.name != 'nt':
            return False, 'non-windows'
        # 1) comtypes path
        com_detail_errors: list[str] = []
        try:
            import comtypes, comtypes.client  # type: ignore
            try:
                # Ensure COM initialized in this thread
                try:
                    comtypes.CoInitialize()
                except Exception:
                    pass
                try:
                    comtypes.client.GetActiveObject('Word.Application')
                    return True, 'comtypes GetActiveObject'
                except Exception as ge:
                    # attempt create
                    try:
                        app = comtypes.client.CreateObject('Word.Application')
                        try:
                            app.Quit()
                        except Exception:
                            pass
                        return True, 'comtypes CreateObject'
                    except Exception as ce:
                        com_detail_errors.append(f'comtypes create failed: {ce}')
                finally:
                    try:
                        comtypes.CoUninitialize()
                    except Exception:
                        pass
            except Exception as inner:
                com_detail_errors.append(f'comtypes op failed: {inner}')
        except Exception as imp_err:
            com_detail_errors.append(f'comtypes import failed: {imp_err}')

        # 2) win32com path
        win_detail_errors: list[str] = []
        try:
            import pythoncom  # type: ignore
            import win32com.client  # type: ignore
            try:
                pythoncom.CoInitialize()
            except Exception:
                pass
            try:
                win32com.client.GetObject(Class='Word.Application')
                return True, 'win32com GetObject'
            except Exception:
                try:
                    wapp = win32com.client.Dispatch('Word.Application')
                    try:
                        wapp.Quit()
                    except Exception:
                        pass
                    return True, 'win32com Dispatch'
                except Exception as d_err:
                    win_detail_errors.append(f'win32com dispatch failed: {d_err}')
            finally:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass
        except Exception as wimp:
            win_detail_errors.append(f'win32com import failed: {wimp}')

        # 3) Registry presence (soft indicator only)
        try:
            if winreg is not None:
                with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, 'Word.Application'):
                    return False, 'Registry ProgID present (Word installed, automation blocked)'
        except Exception:
            pass

        detail = '; '.join(com_detail_errors + win_detail_errors) or 'word not detected'
        return False, detail

    word_com_ok, word_com_detail = _probe_word_com()

    # Embedding model loaded
    embedder_loaded = bool(embedder is not None)

    # LLM model name
    try:
        current_model = llm_get_model_name()
    except Exception:
        current_model = None

    # Dependencies (raw)
    dep_map: dict[str, dict] = {}
    def _add_dep(d):
        dep_map[d['id']] = d
        return d

    _add_dep(_dep('python', True, sys.version))
    _add_dep(_dep('tesseract', bool(TESSERACT_CMD), detail=str(TESSERACT_CMD or 'missing')))
    _add_dep(_dep('poppler', bool(POPPLER_PATH), detail=str(POPPLER_PATH or 'missing'), optional=True))
    _add_dep(_dep('pytesseract', pytesseract is not None, optional=False))
    _add_dep(_dep('pdfminer', pdfminer_extract_text is not None, optional=True))
    _add_dep(_dep('pdfplumber', pdfplumber is not None, optional=True))
    _add_dep(_dep('pypdfium2', True, optional=True))  # imported at top
    _add_dep(_dep('pymupdf', bool(fitz), optional=True))
    _add_dep(_dep('sentence_transformers', 'error' not in info['packages'].get('sentence_transformers', {})))
    _add_dep(_dep('embeddings_db', db_exists, detail=db_path))
    _add_dep(_dep('tts_lib', 'error' not in info['packages'].get('TTS', {}), optional=True))
    _add_dep(_dep('torch', 'error' not in info['packages'].get('torch', {}), optional=True))
    _add_dep(_dep('comtypes', info['automation'].get('comtypes', False)))
    _add_dep(_dep('win32com', info['automation'].get('win32com', False), optional=True))
    _add_dep(_dep('word_com', word_com_ok, detail=word_com_detail, optional=True))
    _add_dep(_dep('cua_adapter', not bool(cua_status.get('error')), detail=str(cua_status.get('available') or cua_status.get('error') or 'unknown'), optional=True))
    _add_dep(_dep('vscode_bin', bool(resolved_vscode), detail=str(resolved_vscode or ''), optional=False))
    _add_dep(_dep('llm_model', bool(current_model), detail=str(current_model or 'unset')))
    _add_dep(_dep('embedder_loaded', embedder_loaded))

    # --- Feature matrix ---
    def feature(name: str, requires: list[str], any_of: list[list[str]] | None = None, optional: bool = False, description: str | None = None):
        """Compute readiness: all 'requires' must be ok; each inner list in any_of requires at least one ok (OR groups)."""
        req_status = {r: dep_map.get(r, {'ok': False}).get('ok', False) for r in requires}
        any_groups = []
        any_groups_ok = True
        if any_of:
            for group in any_of:
                group_states = [(d, dep_map.get(d, {'ok': False}).get('ok', False)) for d in group]
                group_ok = any(state for _, state in group_states)
                any_groups.append({'choices': [d for d, _ in group_states], 'ok': group_ok})
                if not group_ok:
                    any_groups_ok = False
        all_requires_ok = all(req_status.values())
        overall_ok = all_requires_ok and any_groups_ok
        if optional and not overall_ok:
            # Optional features not contributing to global 'ok'
            pass
        return {
            'name': name,
            'ok': overall_ok,
            'requires': req_status,
            'any_of': any_groups or None,
            'optional': optional,
            'description': description,
        }

    features = [
        feature(
            name='OCR Extraction',
            requires=['tesseract', 'pytesseract'],
            any_of=[['poppler', 'pymupdf', 'pypdfium2']],
            description='Image/PDF OCR support (any renderer: PyMuPDF, PDFium, or Poppler).'
        ),
        feature(
            name='PDF Text Extraction (Any)',
            requires=[],
            any_of=[['pdfminer', 'pdfplumber', 'pymupdf', 'pypdfium2']],
            description='At least one rich/light PDF text backend available.'
        ),
        feature(
            name='Embeddings & Semantic Memory',
            requires=['sentence_transformers', 'embeddings_db', 'embedder_loaded'],
            description='Vector embeddings & chat memory persistence.'
        ),
        feature(
            name='LLM Inference',
            requires=['llm_model'],
            any_of=[['torch']],
            description='Core language model response generation.'
        ),
        feature(
            name='Text To Speech',
            requires=['tts_lib'],
            optional=True,
            description='Optional voice synthesis (Coqui TTS).'
        ),
        feature(
            name='Word Automation (Windows)',
            requires=['comtypes'],
            any_of=[['word_com']],
            optional=True,
            description='Open/edit/snap Microsoft Word windows.'
        ),
        feature(
            name='CUA / Snap Assist Automation',
            requires=[],
            any_of=[['cua_adapter']],
            optional=True,
            description='Enhanced window snapping & tile selection.'
        ),
        feature(
            name='VS Code Automation',
            requires=['vscode_bin'],
            description='Open & snap VS Code with target file.'
        ),
        feature(
            name='Browser Pairing for Split Screen',
            requires=['vscode_bin'],
            optional=True,
            description='Relies on runtime detection of browser windows (dynamic).'
        ),
    ]

    info['features'] = features
    # Summarize OCR renderer availability (used to decide if poppler absence is non-critical)
    ocr_renderers = {
        'poppler': bool(POPPLER_PATH),
        'pymupdf': bool(fitz),
        'pypdfium2': True,  # imported at module import; if it failed earlier we'd have noticed
    }
    info['ocr_renderers'] = ocr_renderers
    # Global ok if all non-optional features ok
    critical = [f for f in features if not f.get('optional')]
    info['ok'] = all(f.get('ok') for f in critical)

    # --- Quick readiness flags for UI cards ---
    info['automation']['vscode_ready'] = bool(resolved_vscode)
    info['automation']['word_ready'] = bool(word_com_ok)
    info['automation']['cua_ready'] = bool(cua_status.get('agent_ready') or cua_status.get('available')) and not cua_status.get('error')

    # --- Feature summary (flatten) ---
    features_summary = []
    for f in features:
        missing: list[str] = []
        for dep_id, okv in (f.get('requires') or {}).items():
            if not okv:
                missing.append(dep_id)
        any_groups = f.get('any_of') or []
        for grp in any_groups:
            if not grp.get('ok'):
                # Add group choices as a single OR set descriptor
                choices = grp.get('choices') or []
                if choices:
                    missing.append(' | '.join(choices))
        features_summary.append({
            'name': f['name'],
            'ok': f['ok'],
            'optional': f.get('optional', False),
            'missing': missing,
        })
    info['features_summary'] = features_summary
    missing_critical_raw = {m for fs in features_summary if not fs['optional'] and not fs['ok'] for m in fs['missing']}
    # If OCR feature missing only because of renderer group and at least one non-poppler renderer exists, drop that group from critical
    if any(f['name'] == 'OCR Extraction' and not f['ok'] for f in features_summary):
        # Determine if tesseract core deps are present
        tess_ok = dep_map.get('tesseract', {}).get('ok') and dep_map.get('pytesseract', {}).get('ok')
        if tess_ok and (ocr_renderers.get('pymupdf') or ocr_renderers.get('pypdfium2')):
            # Remove any OR group containing poppler | pymupdf | pypdfium2 from critical if present
            missing_critical_raw = {m for m in missing_critical_raw if 'poppler' not in m and 'pymupdf' not in m and 'pypdfium2' not in m}
    info['missing_critical_dependencies'] = sorted(missing_critical_raw)
    info['missing_any'] = sorted({m for fs in features_summary if not fs['ok'] for m in fs['missing']})
    return info

@app.get('/api/network/laninfo')
def lan_info():
    """Return LAN-accessible URLs for this backend so user can open from another device on same Wi-Fi.

    No external internet required; just devices on same subnet. Frontend can display a copyable list.
    """
    import socket
    hostnames: list[str] = []
    addrs: list[str] = []
    try:
        hn = socket.gethostname()
        hostnames.append(hn)
        try:
            addrs.append(socket.gethostbyname(hn))
        except Exception:
            pass
    except Exception:
        pass
    # Enumerate all interfaces best-effort
    try:
        for info in socket.getaddrinfo(None, 0, family=socket.AF_INET, type=socket.SOCK_STREAM):
            ip = info[4][0]
            if ip and ip not in addrs and not ip.startswith('127.'):
                addrs.append(ip)
    except Exception:
        pass
    # Common dev port assumption 8000
    base_urls = [f"http://{ip}:8000" for ip in addrs]
    # Prune stale clients
    now = time.time()
    stale_cut = now - _LIVE_LAN_CLIENT_TTL
    for ip, ts in list(_LIVE_LAN_CLIENTS.items()):
        if ts < stale_cut:
            _LIVE_LAN_CLIENTS.pop(ip, None)
    # Determine if frontend build is being served
    FE_DIST = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'frontend', 'dist'))
    frontend_served = os.path.isdir(FE_DIST)
    return {
        'hostnames': hostnames,
        'ips': addrs,
        'urls': base_urls if _LIVE_LAN_ENABLED else [],
        'note': 'Open one of these URLs from another device on the same Wi-Fi. Ensure Windows Firewall allows inbound on port 8000.',
        'live_enabled': _LIVE_LAN_ENABLED,
        'connected_clients': sorted(_LIVE_LAN_CLIENTS.keys()),
        'connected_count': len(_LIVE_LAN_CLIENTS),
        'client_ttl_seconds': _LIVE_LAN_CLIENT_TTL,
        'frontend_served': frontend_served,
    }

@app.post('/api/network/live/enable')
def enable_live():
    global _LIVE_LAN_ENABLED
    _LIVE_LAN_ENABLED = True
    return {'live_enabled': True}

@app.post('/api/network/live/disable')
def disable_live():
    global _LIVE_LAN_ENABLED
    _LIVE_LAN_ENABLED = False
    return {'live_enabled': False}

@app.get('/api/network/live/status')
def live_status():
    return {'live_enabled': _LIVE_LAN_ENABLED, 'connected_count': len(_LIVE_LAN_CLIENTS), 'clients': sorted(_LIVE_LAN_CLIENTS.keys())}

@app.get('/api/network/ping')
def live_ping(request: Request):
    if not _LIVE_LAN_ENABLED:
        return {'live_enabled': False}
    ip = request.client.host if request and request.client else 'unknown'
    _LIVE_LAN_CLIENTS[ip] = time.time()
    return {'live_enabled': True, 'your_ip': ip}

# Include Agent Router (Power Mode offline actions)
try:
    from agent_router import router as agent_router
    app.include_router(agent_router)
except Exception as _e:
    logging.warning(f"[Agent] Failed to include agent router: {_e}")

# Include CUA Router (optional; offline-safe placeholder)
try:
    from cua_router import router as cua_router
    app.include_router(cua_router)
except Exception as _e:
    logging.warning(f"[CUA] Failed to include CUA router: {_e}")

# Include Word Router (offline via python-docx)
try:
    from word_router import router as word_router
    app.include_router(word_router)
except Exception as _e:
    logging.warning(f"[Word] Failed to include Word router: {_e}")

# Include Code Router (generic code/text editing)
try:
    from code_router import router as code_router
    app.include_router(code_router)
except Exception as _e:
    logging.warning(f"[Code] Failed to include Code router: {_e}")

# Include Power Mode Router (planning + execute)
try:
    from power_router import router as power_router, OpenDocIntelligentlyRequest, open_doc_intelligently
    app.include_router(power_router)
    
    # Direct endpoint for testing
    @app.post('/api/power/open_doc_intelligently_direct')
    def direct_open_doc(req: OpenDocIntelligentlyRequest):
        print(f"[DIRECT] Received request to open: {req.abs_path}")
        # Idempotency fast-return
        try:
            from power_router import compute_open_doc_signature
            effective_request_id = req.request_id or compute_open_doc_signature(req)
            auto_generated = req.request_id is None
            now = time.time()
            with _RECENT_REQUESTS_LOCK:
                rec = _RECENT_REQUESTS.get(effective_request_id)
                if rec and (now - rec[0] <= _RECENT_REQUESTS_TTL):
                    cached = dict(rec[1])
                    cached['idempotent'] = True
                    cached['idempotent_key'] = effective_request_id
                    if auto_generated:
                        cached['idempotent_auto'] = True
                    print(f"[DIRECT] Result (cached idempotent): {cached}")
                    return cached
        except Exception:
            pass
        result = open_doc_intelligently(req)
        try:
            from power_router import compute_open_doc_signature
            effective_request_id = req.request_id or compute_open_doc_signature(req)
            auto_generated = req.request_id is None
            with _RECENT_REQUESTS_LOCK:
                _RECENT_REQUESTS[effective_request_id] = (time.time(), result)
            result['idempotent_key'] = effective_request_id
            if auto_generated:
                result['idempotent_auto'] = True
        except Exception:
            pass
        print(f"[DIRECT] Result: {result}")
        return result
except Exception as _e:
    logging.warning(f"[Power] Failed to include Power router: {_e}")

# Include Automation Router (merged CUA-only implementation)
try:
    from automation_router import router as automation_router
    app.include_router(automation_router)
except Exception as _e:
    logging.warning(f"[Automation] Failed to include automation router: {_e}")

# Load embedding model
embedder = SentenceTransformer('all-MiniLM-L6-v2')

# Simple in-memory idempotency cache: map request_id -> (timestamp, response_json)
# This is best-effort and resets on process restart. Keep TTL short (e.g., 30s).
_RECENT_REQUESTS = {}
_RECENT_REQUESTS_LOCK = threading.Lock()
_RECENT_REQUESTS_TTL = 30  # seconds

# Configure optional external tools for OCR.
# We auto-detect reasonable defaults so users don't have to set env vars manually.
TESSERACT_CMD = os.environ.get('TESSERACT_CMD')
TESSERACT_LANGS = os.environ.get('TESSERACT_LANGS', 'eng')  # comma or plus separated per tesseract rules
POPPLER_PATH = os.environ.get('POPPLER_PATH')  # path to poppler bin for Windows (poppler/bin or poppler-xx/bin)

_AUTO_TESSERACT_DETECTED = False
_AUTO_POPPLER_DETECTED = False

def _detect_tesseract_on_windows() -> str | None:
    """Try to locate tesseract.exe on Windows via PATH, common install dirs, or registry."""
    # 1) PATH
    try:
        import shutil as _sh
        p = _sh.which('tesseract')
        if p and os.path.isfile(p):
            return p
    except Exception:
        pass
    # 2) Common install locations
    candidates = [
        r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
        r"C:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe",
        os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Programs', 'Tesseract-OCR', 'tesseract.exe'),
        # Common package manager installs
        r"C:\\ProgramData\\chocolatey\\bin\\tesseract.exe",  # Chocolatey
        r"C:\\msys64\\mingw64\\bin\\tesseract.exe",          # MSYS2/MINGW64
        os.path.join(os.environ.get('USERPROFILE', ''), 'scoop', 'apps', 'tesseract', 'current', 'tesseract.exe'),  # Scoop
    ]
    for c in candidates:
        # Expand env vars and ~ if present
        ce = os.path.expanduser(os.path.expandvars(c)) if isinstance(c, str) else c
        if ce and os.path.isfile(ce):
            return ce
    # 2b) Try 'where' as a last resort
    try:
        import subprocess as _subp
        r = _subp.run(['where', 'tesseract'], capture_output=True, text=True, timeout=2)
        if r.returncode == 0:
            for line in (r.stdout or '').splitlines():
                cand = (line or '').strip()
                if cand and os.path.isfile(cand):
                    return cand
    except Exception:
        pass
    # 3) Registry lookup
    try:
        if winreg is not None:
            for hive in (winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER):
                for key_path in (r"SOFTWARE\Tesseract-OCR", r"SOFTWARE\WOW6432Node\Tesseract-OCR"):
                    try:
                        with winreg.OpenKey(hive, key_path) as k:
                            for value_name in ('Path', 'InstallDir', 'TesseractOCRPath'):
                                try:
                                    val, _ = winreg.QueryValueEx(k, value_name)
                                    if val and os.path.isdir(val):
                                        exe = os.path.join(val, 'tesseract.exe')
                                        if os.path.isfile(exe):
                                            return exe
                                except Exception:
                                    continue
                    except FileNotFoundError:
                        continue
    except Exception:
        pass
    return None

def _detect_poppler_on_windows() -> str | None:
    """Try to locate a poppler bin directory (contains pdftoppm / pdftocairo)."""
    if os.name != 'nt':
        return None
    candidates = []
    local = os.environ.get('LOCALAPPDATA','')
    program_files = os.environ.get('ProgramFiles', r"C:\\Program Files")
    program_files_x86 = os.environ.get('ProgramFiles(x86)', r"C:\\Program Files (x86)")
    for root in filter(None, [program_files, program_files_x86, local]):
        try:
            if not os.path.isdir(root):
                continue
            for name in os.listdir(root):
                low = name.lower()
                if 'poppler' in low:
                    base_dir = os.path.join(root, name)
                    # Direct bin
                    cand_direct = os.path.join(base_dir, 'bin')
                    if os.path.isdir(cand_direct) and os.path.isfile(os.path.join(cand_direct, 'pdftoppm.exe')):
                        candidates.append(cand_direct)
                    # Nested Library/bin (zip distribution pattern)
                    cand_lib = os.path.join(base_dir, 'Library', 'bin')
                    if os.path.isdir(cand_lib) and os.path.isfile(os.path.join(cand_lib, 'pdftoppm.exe')):
                        candidates.append(cand_lib)
        except Exception:
            pass
    for c in candidates:
        return c
    # PATH fallback: check if pdftoppm exists
    try:
        import shutil as _sh
        exe = _sh.which('pdftoppm')
        if exe:
            return os.path.dirname(exe)
    except Exception:
        pass
    return None

# Auto-detect if not specified
if os.name == 'nt':
    if not TESSERACT_CMD:
        _detected = _detect_tesseract_on_windows()
        if _detected:
            TESSERACT_CMD = _detected
            _AUTO_TESSERACT_DETECTED = True
    if not POPPLER_PATH:
        _detected_pop = _detect_poppler_on_windows()
        if _detected_pop:
            POPPLER_PATH = _detected_pop
            _AUTO_POPPLER_DETECTED = True

# Simple dependency status endpoint for visibility
@app.get('/api/system/dependencies')
def dependency_status():
    return {
        'tesseract': {
            'configured': bool(TESSERACT_CMD),
            'path': TESSERACT_CMD,
            'auto_detected': _AUTO_TESSERACT_DETECTED,
            'langs': TESSERACT_LANGS,
        },
        'poppler': {
            'configured': bool(POPPLER_PATH),
            'path': POPPLER_PATH,
            'auto_detected': _AUTO_POPPLER_DETECTED,
        },
    }

def _detect_poppler_on_windows() -> str | None:
    """Find Poppler's bin directory containing pdftoppm.exe if installed in common locations."""
    program_files = os.environ.get('ProgramFiles', r"C:\\Program Files")
    program_files_x86 = os.environ.get('ProgramFiles(x86)', r"C:\\Program Files (x86)")
    local_appdata = os.environ.get('LOCALAPPDATA', '')
    choco_install = os.environ.get('ChocolateyInstall', r"C:\\ProgramData\\chocolatey")
    choco_lib = os.path.join(choco_install, 'lib') if choco_install else None
    choco_tools = os.environ.get('ChocolateyToolsLocation', os.path.join(choco_install, 'tools') if choco_install else None)
    bases = [b for b in [
        program_files,
        program_files_x86,
        local_appdata,
        os.path.join(local_appdata, 'Programs') if local_appdata else None,
        choco_lib,
        choco_tools,
    ] if b]
    patterns = [
        os.path.join('*poppler*', 'Library', 'bin'),  # zip installs
        os.path.join('*poppler*', 'bin'),             # generic
        os.path.join('*poppler*', 'tools'),           # Chocolatey
    ]
    found = []
    try:
        for base in bases:
            for pat in patterns:
                for path in glob.glob(os.path.join(base, pat)):
                    if os.path.isfile(os.path.join(path, 'pdftoppm.exe')):
                        found.append(path)
    except Exception:
        pass
    if not found:
        return None
    # Prefer the one with highest version by simple lexicographic sort (good enough)
    found.sort()
    return found[-1]

def _ensure_tesseract_and_poppler() -> bool:
    """(Re)configure Tesseract and Poppler at runtime if needed. Returns True if changed."""
    global TESSERACT_CMD, POPPLER_PATH
    changed = False
    # Tesseract resolution
    try:
        if pytesseract is not None:
            cmd = getattr(pytesseract.pytesseract, 'tesseract_cmd', None)
            # If Windows and cmd is missing or not an absolute file, try to resolve
            if os.name == 'nt':
                import shutil as _sh
                resolved = None
                try:
                    resolved = _sh.which('tesseract')
                except Exception:
                    resolved = None
                if (not cmd) or (not os.path.isabs(cmd)):
                    if resolved and os.path.isfile(resolved):
                        pytesseract.pytesseract.tesseract_cmd = resolved
                        TESSERACT_CMD = resolved
                        changed = True
                # If still not a valid absolute file, probe common locations/registry
                cmd2 = getattr(pytesseract.pytesseract, 'tesseract_cmd', None)
                if (not cmd2) or (os.path.isabs(cmd2) and not os.path.isfile(cmd2)):
                    auto = _detect_tesseract_on_windows()
                    if auto and os.path.isfile(auto):
                        pytesseract.pytesseract.tesseract_cmd = auto
                        TESSERACT_CMD = auto
                        changed = True
    except Exception:
        pass
    # Poppler resolution (Windows only, for pdf2image)
    try:
        if os.name == 'nt' and convert_from_path is not None:
            if (not POPPLER_PATH) or (not os.path.isdir(POPPLER_PATH)):
                auto = _detect_poppler_on_windows()
                if auto and os.path.isdir(auto):
                    POPPLER_PATH = auto
                    changed = True
    except Exception:
        pass
    return changed

# Auto-detect on Windows if not provided
if os.name == 'nt':
    # If env provided but not an absolute path, try to resolve it
    try:
        if TESSERACT_CMD and not os.path.isabs(TESSERACT_CMD):
            import shutil as _sh
            _resolved = _sh.which(TESSERACT_CMD)
            if _resolved:
                TESSERACT_CMD = _resolved
    except Exception:
        pass
    # If env value points to non-existent file, clear it to trigger detection
    if TESSERACT_CMD and not os.path.isfile(TESSERACT_CMD):
        logging.warning(f"[OCR] TESSERACT_CMD does not exist: {TESSERACT_CMD}; attempting auto-detect.")
        TESSERACT_CMD = None
    if not TESSERACT_CMD and pytesseract is not None:
        auto_tess = _detect_tesseract_on_windows()
        if auto_tess:
            TESSERACT_CMD = auto_tess
            logging.info(f"[OCR] Auto-detected Tesseract: {TESSERACT_CMD}")
        else:
            logging.warning("[OCR] Could not auto-detect Tesseract. Install it or set TESSERACT_CMD.")
    if not POPPLER_PATH and convert_from_path is not None:
        auto_poppler = _detect_poppler_on_windows()
        if auto_poppler:
            POPPLER_PATH = auto_poppler
            logging.info(f"[OCR] Auto-detected Poppler bin: {POPPLER_PATH}")
        else:
            logging.info("[OCR] Poppler not found; will use PDFium fallback if available.")

# Apply Tesseract CMD to pytesseract if available
if pytesseract is not None and TESSERACT_CMD:
    try:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        logging.info(f"[OCR] Using Tesseract at {TESSERACT_CMD}")
    except Exception as _e:
        logging.warning(f"[OCR] Failed to set Tesseract cmd: {_e}")

# Final validation: if version cannot be obtained, try one last auto-detect on Windows
if pytesseract is not None:
    try:
        _ver = str(pytesseract.get_tesseract_version())
        if not _ver and os.name == 'nt':
            # attempt detect & set again
            _auto = _detect_tesseract_on_windows()
            if _auto and os.path.isfile(_auto):
                pytesseract.pytesseract.tesseract_cmd = _auto
                TESSERACT_CMD = _auto
                logging.info(f"[OCR] Recovered Tesseract path: {TESSERACT_CMD}")
    except Exception:
        if os.name == 'nt':
            _auto = _detect_tesseract_on_windows()
            if _auto and os.path.isfile(_auto):
                try:
                    pytesseract.pytesseract.tesseract_cmd = _auto
                    TESSERACT_CMD = _auto
                    logging.info(f"[OCR] Recovered Tesseract path: {TESSERACT_CMD}")
                except Exception:
                    pass

# -------------------- OCR/Text Extraction Helpers --------------------
def _extract_text_pdf_best_effort(path: str) -> str:
    """Try multiple libraries for text-based PDFs before OCR: PyMuPDF -> pdfplumber -> pdfminer.six -> PyPDF2."""
    # 1) PyMuPDF
    if fitz is not None:
        try:
            with fitz.open(path) as doc:
                parts = []
                for page in doc:
                    try:
                        t = page.get_text() or ''
                    except Exception:
                        t = ''
                    if t:
                        parts.append(t)
                combined = "\n".join(parts)
                if combined.strip():
                    return combined
        except Exception:
            pass
    # 2) pdfplumber
    if pdfplumber is not None:
        try:
            text = ''
            with pdfplumber.open(path) as pdf:
                for pg in pdf.pages:
                    try:
                        text += (pg.extract_text() or '') + '\n'
                    except Exception:
                        continue
            if text.strip():
                return text
        except Exception:
            pass
    # 3) pdfminer.six (pure Python)
    if pdfminer_extract_text is not None:
        try:
            text = pdfminer_extract_text(path) or ''
            if text.strip():
                return text
        except Exception:
            pass
    # 4) PyPDF2
    try:
        reader = PyPDF2.PdfReader(path)
        pages = [p.extract_text() or '' for p in reader.pages]
        text = '\n'.join(pages)
        if text.strip():
            return text
    except Exception:
        pass
    return ''

def resolve_tags_to_context(tags: list[str]) -> Optional[str]:
    """Best-effort: given a list of tag strings, return a short joined preview string
    by looking up mem_folder and mem_item tables. This mirrors frontend resolveTagsToContext
    so regenerate requests get the same context even if the client didn't send mem_context.
    """
    if not tags:
        return None
    try:
        conn = sqlite3.connect(DB_PATH)
        c = conn.cursor()
        # Load folders for mapping tag/name -> folder
        c.execute('SELECT id, name, tag FROM mem_folder')
        folders = c.fetchall() or []
        folder_by_key = {}
        for fid, name, tag in folders:
            if tag:
                folder_by_key[str(tag).strip().lower()] = {'id': fid, 'name': name}
            if name:
                folder_by_key[str(name).strip().lower()] = {'id': fid, 'name': name}

        contexts = []
        for t in tags:
            key = str(t).lower()
            if key in folder_by_key:
                folder = folder_by_key[key]
                try:
                    c.execute('SELECT substr(text_content,1,500) as preview, title, filename FROM mem_item WHERE folder_id=? ORDER BY id DESC LIMIT 5', (folder['id'],))
                    rows = c.fetchall() or []
                    previews = [r[0] or r[1] or r[2] or '' for r in rows]
                    previews = [p for p in previews if p]
                    if previews:
                        contexts.append(f"Folder: {folder['name']}\n" + '\n---\n'.join(previews))
                except Exception:
                    continue
            else:
                try:
                    like = f"%{t}%"
                    c.execute('SELECT substr(text_content,1,500), filename FROM mem_item WHERE tags LIKE ? OR filename LIKE ? ORDER BY id DESC LIMIT 5', (like, like))
                    rows = c.fetchall() or []
                    previews = [(r[0] or r[1] or '') for r in rows]
                    previews = [p for p in previews if p]
                    if previews:
                        contexts.append(f"Tag: #{t}\n" + '\n---\n'.join(previews))
                except Exception:
                    continue
        conn.close()
        if not contexts:
            return None
        joined = '\n====\n'.join(contexts)
        # Trim to a reasonable length
        return joined[:3000]
    except Exception:
        try:
            conn.close()
        except Exception:
            pass
        return None

def _ocr_pdf_via_fitz(path: str, lang: str) -> str:
    """Render each page using PyMuPDF and OCR with Tesseract; avoids Poppler dependency."""
    if fitz is None or pytesseract is None:
        return ''
    out = ''
    try:
        with fitz.open(path) as doc:
            mat = fitz.Matrix(2.0, 2.0)  # ~200 DPI
            for page in doc:
                try:
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img_bytes = pix.tobytes("png")
                    pil = Image.open(io.BytesIO(img_bytes))
                    out += pytesseract.image_to_string(pil, lang=lang) + '\n'
                except Exception:
                    continue
    except Exception:
        return ''
    return out

# SQLite DB for chat embeddings
DB_PATH = 'chat_embeddings.db'
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS chat (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        chat_id TEXT,
        user TEXT,
        llm TEXT,
        embedding BLOB,
        timestamp INTEGER,
        tags TEXT,
        doc_info TEXT,
        service TEXT,
        chat_name TEXT
    )''')
    # Memory manager tables
    c.execute('''CREATE TABLE IF NOT EXISTS mem_folder (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT,
        tag TEXT,
        created_at INTEGER
    )''')
    c.execute('''CREATE TABLE IF NOT EXISTS mem_item (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        folder_id INTEGER,
        title TEXT,
        filename TEXT,
        text_content TEXT,
        tags TEXT,
        embedding BLOB,
        created_at INTEGER,
        FOREIGN KEY(folder_id) REFERENCES mem_folder(id)
    )''')
    conn.commit()
    # Track artifacts generated by power mode for cleanup
    c.execute('''CREATE TABLE IF NOT EXISTS chat_artifact (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        chat_id TEXT,
        service TEXT,
        path TEXT,
        created_at INTEGER
    )''')
    conn.commit()
    # Ensure older databases get the new 'service' column if missing
    try:
        c.execute("PRAGMA table_info(chat)")
        cols = [r[1] for r in c.fetchall()]
        if 'service' not in cols:
            c.execute("ALTER TABLE chat ADD COLUMN service TEXT DEFAULT 'default'")
            conn.commit()
        if 'chat_name' not in cols:
            # Add chat_name column and default to NULL for existing rows
            c.execute("ALTER TABLE chat ADD COLUMN chat_name TEXT DEFAULT NULL")
            conn.commit()
        if 'tags' not in cols:
            # Add tags column to persist mem_tags for each chat row (optional)
            c.execute("ALTER TABLE chat ADD COLUMN tags TEXT DEFAULT NULL")
            conn.commit()
        # Create chat_state table if missing to persist per-chat Power Mode state
        c.execute('''CREATE TABLE IF NOT EXISTS chat_state (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chat_id TEXT,
            service TEXT,
            persistent_tags TEXT,
            doc_path TEXT,
            created_at INTEGER,
            updated_at INTEGER
        )''')
        conn.commit()
    except Exception as e:
        logging.debug(f"Could not migrate chat table to add 'service' column: {e}")
    conn.close()
init_db()
print("Database initialized.")

class ChatRequest(BaseModel):
    chat_id: str
    text: str
    file: Optional[str] = None
    doc_info: Optional[str] = None
    service: Optional[str] = 'default'
    mem_context: Optional[str] = None
    # Optional list of tag strings selected by the user (e.g. ['projectA', 'notes'])
    mem_tags: Optional[List[str]] = None
    # If True, replace the last AI response for this chat (same chat_id) instead of inserting a new row
    replace_last: Optional[bool] = False
    # Optional: replace a specific DB row id (preferred). If provided, server will update that row's llm.
    replace_id: Optional[int] = None
    # Optional idempotency key supplied by client. If present, server will return cached result
    # for previously-seen request_id values (short TTL) to avoid duplicate processing.
    request_id: Optional[str] = None

# Translation endpoint using LLM and TTS (after app is defined)
class TranslateRequest(BaseModel):
    text: str
    from_lang: str
    to_lang: str

@app.post('/api/translate')
async def translate_endpoint(request: Request):
    data = await request.json()
    text = data.get('text', '')
    from_lang = data.get('from') or data.get('from_lang') or 'en'
    to_lang = data.get('to') or data.get('to_lang') or 'en'
    # Use LLM to translate with professional translator system prompt
    system_prompt = f"""
You are a professional translator AI.  
Your task is to translate text accurately, naturally, and in context.  

Rules:
1. Always translate from {from_lang} → {to_lang}.  
2. Preserve meaning, tone, and style of the original text.  
3. Do not add or omit information.  
4. Keep placeholders, numbers, and formatting unchanged.  
5. If the text is ambiguous, give the most natural interpretation.  
6. Output only the translated text — no explanations unless asked.  
"""
    prompt = system_prompt + f"\nText: {text}"
    translated = generate_response(prompt)
    # Do NOT generate TTS here. Only return text.
    return {"text": translated}

# -------------------- Ollama Model Management --------------------
import subprocess
import platform
import requests as _req

RECOMMENDED_MODELS = [
    {"name": "llama3.1:8b", "estimated_size": "4-5 GB"},
    {"name": "qwen2.5:7b", "estimated_size": "4-5 GB"},
    {"name": "mistral:7b", "estimated_size": "4-5 GB"},
    {"name": "phi3:mini", "estimated_size": "1.8 GB"},
]

@app.get('/api/ollama/models')
def list_ollama_models():
    """Return installed Ollama models and a recommended list, plus current model used by backend."""
    installed = []
    try:
        r = _req.get('http://127.0.0.1:11434/api/tags', timeout=3)
        if r.ok:
            data = r.json() or {}
            for m in data.get('models', []):
                installed.append({
                    'name': m.get('name'),
                    'size': m.get('size'),
                    'digest': m.get('digest'),
                    'modified_at': m.get('modified_at')
                })
    except Exception as e:
        logging.warning(f"[Ollama] Could not fetch installed models: {e}")
    return {
        'current_model': llm_get_model_name(),
        'installed': installed,
        'recommended': RECOMMENDED_MODELS
    }

class InstallModelRequest(BaseModel):
    model: str

@app.post('/api/ollama/install')
def install_ollama_model(req: InstallModelRequest):
    """Run 'ollama pull <model>' on the host machine. Returns success flag and output snippet."""
    model = (req.model or '').strip()
    if not model:
        return JSONResponse(status_code=400, content={'error': 'model is required'})
    try:
        # Use shell to ensure PATH resolution; on Windows, PowerShell is fine to wrap the command
        if platform.system().lower().startswith('win'):
            cmd = ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", f"ollama pull {model}"]
        else:
            cmd = ["/bin/sh", "-lc", f"ollama pull {model}"]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=60*30)
        success = proc.returncode == 0
        out = (proc.stdout or '')[-4000:]
        err = (proc.stderr or '')[-4000:]
        return { 'success': success, 'stdout': out, 'stderr': err }
    except subprocess.TimeoutExpired:
        return JSONResponse(status_code=504, content={'error': 'Installation timed out'})
    except Exception as e:
        logging.exception('Ollama install failed')
        return JSONResponse(status_code=500, content={'error': str(e)})

class ModelSetRequest(BaseModel):
    model: str

@app.get('/api/model')
def get_model():
    return { 'model': llm_get_model_name() }

@app.post('/api/model')
def set_model(req: ModelSetRequest):
    name = (req.model or '').strip()
    if not name:
        return JSONResponse(status_code=400, content={'error': 'model is required'})
    try:
        llm_set_model_name(name)
        return { 'ok': True, 'model': name }
    except Exception as e:
        return JSONResponse(status_code=500, content={'error': str(e)})

@app.post('/api/chat')
def chat_endpoint(req: ChatRequest):
    import numpy as np
    # Quick idempotency check: if client provided request_id and we have a recent cached response, return it
    try:
        rid_key = getattr(req, 'request_id', None)
    except Exception:
        rid_key = None
    if rid_key:
        now = time.time()
        with _RECENT_REQUESTS_LOCK:
            entry = _RECENT_REQUESTS.get(rid_key)
            if entry:
                ts, resp = entry
                if now - ts <= _RECENT_REQUESTS_TTL:
                    # Return cached response (assumed to be small JSON-serializable dict)
                    return resp
                else:
                    # Stale entry: remove
                    try:
                        del _RECENT_REQUESTS[rid_key]
                    except Exception:
                        pass

    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()

    # --- Include previous N interactions from this chat_id as conversational context ---
    try:
        c.execute('SELECT user, llm, doc_info FROM chat WHERE chat_id=? ORDER BY id DESC LIMIT 3', (req.chat_id,))
        prev_rows = c.fetchall()
    except Exception:
        prev_rows = []
    # Build previous conversation text (oldest first)
    prev_context = ''
    if prev_rows:
        try:
            # If client requested replace_last, avoid including the last AI reply
            # so the LLM doesn't simply repeat the previous answer when regenerating.
            if getattr(req, 'replace_last', False) and len(prev_rows) > 0:
                # prev_rows is ordered DESC (most recent first) — clear llm of most recent row
                try:
                    # convert to list to mutate
                    pr = list(prev_rows)
                    first_user, first_llm, first_di = pr[0]
                    pr[0] = (first_user, None, first_di)
                    rows_to_iterate = list(reversed(pr))
                except Exception:
                    rows_to_iterate = reversed(prev_rows)
            else:
                rows_to_iterate = reversed(prev_rows)
            for ur, lr, di in rows_to_iterate:
                if ur:
                    prev_context += f"User: {ur}\n"
                if lr:
                    prev_context += f"AI: {lr}\n"
        except Exception:
            prev_context = ''

    # --- Embedding-based similarity lookup: only run for Research service ---
    svc = (req.service or '').lower() if req.service else ''
    is_research = ('research' in svc) or (svc == 'research_assistant')
    match_info = []
    if is_research:
        # compute embedding for the input and compare to stored user embeddings from the same chat_id/service
        embedding = embedder.encode(req.text)
        if req.service:
            c.execute('SELECT user, embedding, doc_info FROM chat WHERE chat_id=? AND service=? AND user IS NOT NULL AND embedding IS NOT NULL', (req.chat_id, req.service))
        else:
            c.execute('SELECT user, embedding, doc_info FROM chat WHERE chat_id=? AND user IS NOT NULL AND embedding IS NOT NULL', (req.chat_id,))
        rows = c.fetchall()
        # Dynamic threshold based on user input length
        input_len = len(req.text.strip().split())
        if input_len < 10:
            threshold = 0.35
        elif input_len < 30:
            threshold = 0.5
        elif input_len < 80:
            threshold = 0.6
        else:
            threshold = 0.7
        for user_text, emb_bytes, doc_info in rows:
            if not user_text or len(user_text.strip()) < 3 or not emb_bytes:
                continue  # skip trivial/empty texts
            emb = np.frombuffer(emb_bytes, dtype=np.float32)
            # cosine similarity
            try:
                sim = float(np.dot(embedding, emb) / (np.linalg.norm(embedding) * np.linalg.norm(emb)))
            except Exception:
                continue
            sim = abs(sim)
            if sim >= threshold:
                match_info.append((user_text, doc_info, sim))

    # Build service-specific system prompt based on req.service
    if 'beginner' in svc or 'teacher' in svc:
        system_prompt = """
You are a Beginner-Friendly Teacher AI. Explain topics step by step in clear, simple language.  

Rules:
1. Use short sentences, avoid jargon; define terms simply.  
2. Always show numbered steps for solutions.  
3. For math: compute step by step, show intermediate results.  
4. For coding: provide runnable examples + short explanation + expected output.  
5. For concepts: start with 1-line summary → steps/ideas → one simple example.  
6. End with a quick check question + 1 practice task.  
7. Be encouraging, patient, never condescending.  
8. Only state facts you know or given by user; if unsure, say so.  
9. Adapt depth to user’s level if known.  
10. Use clean formatting: headings, steps, bullets, ✅ for confirmations, ⚠️ for warnings.  

Goal: Act like a supportive teacher who makes learning easy, clear, and step-by-step for beginners.
"""
    elif 'daily' in svc or 'companion' in svc:
        system_prompt = """
You are a Daily Chat Companion.  
Your goal is to keep conversations light, short, and friendly.  
Give crisp answers in simple, everyday language.  
Avoid long explanations unless the user asks.  
Focus on being clear, approachable, and easy to understand.
"""
    elif 'homework' in svc:
        system_prompt = """
You are a Homework Helper.  
Your job is to assist students by giving factual, accurate answers.  
Always make sure responses are aligned with established theories, textbooks, or widely accepted academic knowledge.  
Do not speculate or invent facts.  
Explain step by step when needed, and keep the tone supportive and clear.  
If the information is uncertain, clearly state that.
"""
    else:
        system_prompt = """
You are a Research Partner AI. Your role is to assist the user with thoughtful, accurate, and well-structured responses.

Guidelines:
1. **Tone & Personality**
- Be collaborative, like a supportive research partner or co-thinker.
- Maintain professionalism but allow light empathy and encouragement (e.g., "That’s a great insight!" or "We can refine this further.").
- Avoid being overly casual or overly robotic.

2. **Factual Integrity**
- Only make claims that are factual, verifiable, or explicitly provided by the user.
- If uncertain, acknowledge limits instead of guessing.
- When using external knowledge, clearly distinguish between well-established facts and interpretations.

3. **Depth & Rigor**
- Provide structured, detailed explanations.
- Support reasoning with definitions, examples, or simple abstraction when useful.
- Suggest alternative perspectives or methods if relevant.

4. **Clarity & Usefulness**
- Present information clearly, concisely, and logically ordered.
- Break down complex concepts into understandable steps without oversimplifying.
- Where helpful, summarize key takeaways at the end.

5. **Boundaries**
- Never invent sources or false references.
- Do not speculate beyond what the user has given, unless clearly marked as hypothesis.
- Respect user instructions and constraints fully.

Overall, act as a reliable, insightful, and empathetic partner in research, analysis, and problem-solving.
"""

    # Append embedding-based matches (only present when is_research)
    if match_info:
        system_prompt += "\nBelow is info from previous chat messages and docs that may help you:"
        for user_text, doc_info, sim in match_info:
            system_prompt += f"\n- Info: {user_text}"
            if doc_info:
                system_prompt += f" (Document: {doc_info})"
            system_prompt += f" [similarity: {sim:.2f}]"

    # Append the previous conversation context (last 3 interactions) so LLM can consider recent history
    if prev_context:
        system_prompt += "\nPrevious conversation context:\n" + prev_context
    # Log matches and prompt
    print(f"[DEBUG] Matching embeddings for chat_id {req.chat_id}: {match_info}")
    print(f"[DEBUG] System prompt sent to LLM:\n{system_prompt}\nUser input: {req.text}")

    # Final prompt to LLM
    # If the client passed memory context (from #tag), include it as supplemental context.
    # Otherwise attempt to extract inline #tags from the user text and resolve them from the DB.
    # Priority for mem_context sources:
    # 1) client-provided mem_context (explicit string)
    # 2) client-provided mem_tags (list of tag strings)
    # 3) request.tags (list or comma-separated string)
    # 4) inline tags found in the text (e.g. #tag)
    # 5) most recent saved tags in DB for this chat_id (and service, if provided)
    mem_ctx = getattr(req, 'mem_context', None)
    mem_ctx_source = None
    used_tags: list[str] = []  # track which tags we actually used to build context
    if not mem_ctx:
        # If client provided mem_tags (selected chips), use them first
        try:
            if getattr(req, 'mem_tags', None):
                tags_list = []
                for t in req.mem_tags:
                    if not t:
                        continue
                    s = str(t).strip()
                    if not s:
                        continue
                    if s.startswith('#'):
                        s = s[1:]
                    if s:
                        tags_list.append(s)
                if tags_list:
                    mem_ctx = resolve_tags_to_context(tags_list)
                    if mem_ctx:
                        mem_ctx_source = 'mem_tags'
                        used_tags = list(tags_list)
        except Exception:
            mem_ctx = None
    # Fallback 2: request.tags (string or array)
    if not mem_ctx:
        try:
            if hasattr(req, 'tags'):
                rt = getattr(req, 'tags')
                tags_list = []
                if isinstance(rt, list):
                    for t in rt:
                        if not t:
                            continue
                        s = str(t).strip()
                        if not s:
                            continue
                        if s.startswith('#'):
                            s = s[1:]
                        if s:
                            tags_list.append(s)
                elif isinstance(rt, str):
                    parts = [p.strip() for p in rt.split(',')]
                    for p in parts:
                        if not p:
                            continue
                        s = p[1:] if p.startswith('#') else p
                        if s:
                            tags_list.append(s)
                if tags_list:
                    mem_ctx = resolve_tags_to_context(tags_list)
                    if mem_ctx:
                        mem_ctx_source = 'req.tags'
                        used_tags = list(tags_list)
        except Exception:
            mem_ctx = None
    if not mem_ctx:
        try:
            # find inline tags like #tag or #tag_name
            inline_tags = re.findall(r"#([a-zA-Z0-9_-]+)", req.text or "")
            # keep unique, preserve order
            if inline_tags:
                seen = {}
                uniq = []
                for t in inline_tags:
                    k = str(t).strip()
                    if not k: continue
                    lk = k.lower()
                    if lk not in seen:
                        seen[lk] = True
                        uniq.append(k)
                if uniq:
                    mem_ctx = resolve_tags_to_context(uniq)
                    if mem_ctx:
                        mem_ctx_source = 'inline'
                        used_tags = list(uniq)
        except Exception:
            mem_ctx = None
    # Fallback 4: look up most recent saved tags for this chat in DB
    if not mem_ctx:
        try:
            if req.service:
                c.execute('SELECT tags FROM chat WHERE chat_id=? AND service=? AND tags IS NOT NULL AND TRIM(tags)<>"" ORDER BY id DESC LIMIT 1', (req.chat_id, req.service))
            else:
                c.execute('SELECT tags FROM chat WHERE chat_id=? AND tags IS NOT NULL AND TRIM(tags)<>"" ORDER BY id DESC LIMIT 1', (req.chat_id,))
            row = c.fetchone()
            if row and row[0]:
                tags_str = str(row[0])
                tags_list = []
                for p in tags_str.split(','):
                    s = p.strip()
                    if not s:
                        continue
                    if s.startswith('#'):
                        s = s[1:]
                    if s:
                        tags_list.append(s)
                if tags_list:
                    mem_ctx = resolve_tags_to_context(tags_list)
                    if mem_ctx:
                        mem_ctx_source = 'db.tags'
                        used_tags = list(tags_list)
        except Exception:
            pass

    # Log whether memory context came from client or was server-resolved (helps testing)
    try:
        if getattr(req, 'mem_context', None):
            logging.info(f"[CHAT] request_id={getattr(req,'request_id',None)} mem_context=client-provided len={len(str(getattr(req,'mem_context') or ''))}")
        elif mem_ctx:
            logging.info(f"[CHAT] request_id={getattr(req,'request_id',None)} mem_context=resolved source={mem_ctx_source} len={len(str(mem_ctx))}")
        else:
            logging.info(f"[CHAT] request_id={getattr(req,'request_id',None)} mem_context=none")
    except Exception:
        pass

    if mem_ctx:
        system_prompt += "\nMemory context from user request:\n" + mem_ctx

    # If we have effective tags, surface any explicit URLs from matching memory items so the model can't miss links
    try:
        if used_tags:
            # Gather previews containing URLs for each tag; keep results small and deduped
            conn2 = sqlite3.connect(DB_PATH)
            c2 = conn2.cursor()
            previews = []
            seen_ids = set()
            for t in used_tags:
                like = f"%{t}%"
                try:
                    c2.execute(
                        'SELECT id, substr(text_content,1,1000) as preview, filename FROM mem_item '\
                        'WHERE (tags LIKE ? OR filename LIKE ?) AND ('
                        'text_content LIKE "%http%" OR text_content LIKE "%www.%") '
                        'ORDER BY id DESC LIMIT 20', (like, like)
                    )
                    for rid, prev, fn in c2.fetchall() or []:
                        if rid in seen_ids:
                            continue
                        seen_ids.add(rid)
                        previews.append(prev or '')
                except Exception:
                    continue
            try:
                conn2.close()
            except Exception:
                pass
            # Extract URL strings
            urls = []
            if previews:
                url_re = re.compile(r'https?://[^\s)]+', re.IGNORECASE)
                for pv in previews:
                    for m in url_re.findall(pv or ''):
                        u = m.strip().rstrip('.,);]')
                        if u and u not in urls:
                            urls.append(u)
                        if len(urls) >= 20:
                            break
                    if len(urls) >= 20:
                        break
            if urls:
                # Append a compact section to the prompt so LLM can reference them explicitly
                system_prompt += "\nRelevant links from memory (by tags: " + ", ".join(used_tags) + "):\n" + "\n".join([f"- {u}" for u in urls])
    except Exception:
        pass

    # Also include doc_info (filenames / doc tags) if present to help RAG
    if getattr(req, 'doc_info', None):
        try:
            system_prompt += "\nDocument info:\n" + str(req.doc_info)
        except Exception:
            pass

    final_prompt = system_prompt + "\nUser: " + req.text
    ai_text = generate_response(final_prompt)

    # Only compute/store embeddings for research chats
    emb_to_store = None
    if is_research:
        try:
            emb_to_store = embedder.encode(req.text)
        except Exception:
            emb_to_store = None

    # Save user and LLM output in one row; embedding stored only for research (or None otherwise)
    emb_blob = emb_to_store.tobytes() if emb_to_store is not None else None

    # Determine chat_name: if this chat_id already has a chat_name, reuse it; otherwise create default 'New Chat N'
    c.execute('SELECT DISTINCT chat_name FROM chat WHERE chat_id=? AND chat_name IS NOT NULL', (req.chat_id,))
    existing_name_row = c.fetchone()
    if existing_name_row and existing_name_row[0]:
        chat_name_to_use = existing_name_row[0]
    else:
        # compute next New Chat number globally
        c.execute("SELECT DISTINCT chat_name FROM chat WHERE chat_name LIKE 'New Chat %'")
        existing = [r[0] for r in c.fetchall() if r and r[0]]
        max_n = 0
        for en in existing:
            try:
                n = int(en.replace('New Chat ', '').strip())
                if n > max_n:
                    max_n = n
            except Exception:
                continue
        chat_name_to_use = f'New Chat {max_n + 1}'

    # Normalize incoming tags into a comma-separated string for storage (if provided).
    # Accept either `tags` or `mem_tags` from different frontends.
    incoming_tags_str = None
    try:
        tag_source = None
        if getattr(req, 'mem_tags', None):
            tag_source = getattr(req, 'mem_tags')
        elif getattr(req, 'tags', None):
            tag_source = getattr(req, 'tags')
        if tag_source:
            incoming_tags = [str(t).strip() for t in tag_source if t and str(t).strip()]
            if incoming_tags:
                incoming_tags_str = ','.join(incoming_tags)
    except Exception:
        incoming_tags_str = None

    if getattr(req, 'replace_last', False) or getattr(req, 'replace_id', None):
        try:
            # If client provided an explicit replace_id, prefer that (precise replace-by-id)
            if getattr(req, 'replace_id', None):
                try:
                    rid = int(req.replace_id)
                    c.execute('SELECT id FROM chat WHERE id=?', (rid,))
                    if c.fetchone():
                        # If client provided tags, update tags; otherwise preserve existing tags
                        if incoming_tags_str is not None:
                            c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=?, tags=? WHERE id=?',
                                      (ai_text, req.doc_info, req.service, incoming_tags_str, rid))
                        else:
                            c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=? WHERE id=?',
                                      (ai_text, req.doc_info, req.service, rid))
                        conn.commit()
                        updated = True
                    else:
                        updated = False
                except Exception:
                    updated = False
            else:
                updated = False

            # If not updated by id, fall back to best-effort replace_last behavior
            if not updated and getattr(req, 'replace_last', False):
                # Prefer to update a row that matches the same user text (most precise)
                c.execute('SELECT id FROM chat WHERE chat_id=? AND user=? ORDER BY id DESC LIMIT 1', (req.chat_id, req.text))
                row = c.fetchone()
                if row:
                    row_id = row[0]
                    # Preserve existing tags if client did not send tags
                    if incoming_tags_str is not None:
                        c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=?, tags=? WHERE id=?',
                                  (ai_text, req.doc_info, req.service, incoming_tags_str, row_id))
                    else:
                        c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=? WHERE id=?',
                                  (ai_text, req.doc_info, req.service, row_id))
                    conn.commit()
                else:
                    # Fallback: update the last AI row for this chat_id
                    c.execute('SELECT id FROM chat WHERE chat_id=? AND llm IS NOT NULL ORDER BY id DESC LIMIT 1', (req.chat_id,))
                    last = c.fetchone()
                    if last:
                        last_id = last[0]
                        if incoming_tags_str is not None:
                            c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=?, tags=? WHERE id=?',
                                      (ai_text, req.doc_info, req.service, incoming_tags_str, last_id))
                        else:
                            c.execute('UPDATE chat SET llm=?, timestamp=strftime("%s","now"), doc_info=?, service=? WHERE id=?',
                                      (ai_text, req.doc_info, req.service, last_id))
                        conn.commit()
                    else:
                        # If nothing to replace, insert as new
                        c.execute('INSERT INTO chat (chat_id, user, llm, embedding, timestamp, doc_info, service, chat_name, tags) VALUES (?, ?, ?, ?, strftime("%s", "now"), ?, ?, ?, ?)',
                                  (req.chat_id, req.text, ai_text, emb_blob, req.doc_info, req.service, chat_name_to_use, incoming_tags_str))
                        conn.commit()
        except Exception:
            # On error, fall back to inserting new row
            c.execute('INSERT INTO chat (chat_id, user, llm, embedding, timestamp, doc_info, service, chat_name, tags) VALUES (?, ?, ?, ?, strftime("%s", "now"), ?, ?, ?, ?)',
                      (req.chat_id, req.text, ai_text, emb_blob, req.doc_info, req.service, chat_name_to_use, incoming_tags_str))
            conn.commit()
    else:
        # Insert new row; include tags if provided
        c.execute('INSERT INTO chat (chat_id, user, llm, embedding, timestamp, doc_info, service, chat_name, tags) VALUES (?, ?, ?, ?, strftime("%s", "now"), ?, ?, ?, ?)',
                  (req.chat_id, req.text, ai_text, emb_blob, req.doc_info, req.service, chat_name_to_use, incoming_tags_str))
        conn.commit()
    # Update chat_state persistent tags if any were provided this turn
    try:
        if incoming_tags_str is not None:
            tag_list = [t.strip() for t in incoming_tags_str.split(',') if t and t.strip()]
        else:
            tag_list = None
    except Exception:
        tag_list = None
    try:
        conn2 = sqlite3.connect(DB_PATH)
        _upsert_chat_state(conn2, req.chat_id, req.service, tag_list, None)
        conn2.close()
    except Exception:
        pass
    conn.close()
    resp = {"text": ai_text, "isVoiceMessage": False, "voiceUrl": None}
    # Store in idempotency cache if request_id provided
    try:
        if getattr(req, 'request_id', None):
            now = time.time()
            with _RECENT_REQUESTS_LOCK:
                # Keep a simple copy; entries will be pruned by TTL on access
                _RECENT_REQUESTS[req.request_id] = (now, resp)
    except Exception:
        pass
    return resp
# Endpoint to rename a chat (update all rows for a chat_id with new name)
from fastapi import Body
class RenameChatRequest(BaseModel):
    chat_id: str
    new_name: str

@app.post('/api/rename_chat')
def rename_chat(req: RenameChatRequest):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    # Update chat_name for all rows for this chat_id
    c.execute('UPDATE chat SET chat_name = ? WHERE chat_id = ?', (req.new_name, req.chat_id))
    conn.commit()
    conn.close()
    return {"status": "ok", "chat_id": req.chat_id, "new_name": req.new_name}

@app.post('/api/upload')
def upload_file(file: UploadFile = File(...)):
    """Upload a file and, if it's an image, attempt OCR so frontend can append text for LLM context.
    Returns: { url, name, ocr_text? }
    """
    upload_dir = os.path.join(os.path.dirname(__file__), 'uploads')
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, file.filename)
    raw = file.file.read()
    with open(file_path, 'wb') as f:
        f.write(raw)
    ocr_text = None
    try:
        ext = file.filename.split('.')[-1].lower()
        if ext in ('png','jpg','jpeg','bmp','tif','tiff') and 'ocr' not in (file.filename.lower()):
            if 'pytesseract' in globals() and pytesseract is not None:
                try:
                    img = Image.open(io.BytesIO(raw))
                    ocr_text = pytesseract.image_to_string(img, lang=TESSERACT_LANGS)
                    if ocr_text:
                        # Basic cleanup
                        ocr_text = ocr_text.replace('\r','').strip()
                        # Discard if it's just whitespace or too short
                        if len(ocr_text.strip()) < 3:
                            ocr_text = None
                except Exception as _e:
                    logging.warning(f"[OCR] image OCR failed for {file.filename}: {_e}")
            else:
                logging.info("[OCR] pytesseract not available; skipping image OCR")
    except Exception as _e_all:
        logging.warning(f"[OCR] unexpected error while processing {file.filename}: {_e_all}")
    return {"url": file_path, "name": file.filename, **({"ocr_text": ocr_text} if ocr_text else {})}

# -------------------- Memory Manager Endpoints --------------------
@app.get('/api/memory/folders')
def list_mem_folders():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('SELECT id, name, tag, created_at FROM mem_folder ORDER BY id DESC')
    rows = c.fetchall()
    conn.close()
    folders = [{'id': r[0], 'name': r[1], 'tag': r[2], 'created_at': r[3]} for r in rows]
    return {'folders': folders}

@app.post('/api/memory/folders')
async def create_mem_folder(req: Request):
    # Expect JSON body: { name: str, tag?: str }
    try:
        data = await req.json()
    except Exception:
        data = {}
    name = data.get('name') if isinstance(data, dict) else None
    tag = data.get('tag') if isinstance(data, dict) else None
    if not name:
        return JSONResponse({'error': 'name required'}, status_code=400)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('INSERT INTO mem_folder (name, tag, created_at) VALUES (?, ?, strftime("%s", "now"))', (name, tag))
    conn.commit()
    conn.close()
    return {'status': 'ok'}

@app.get('/api/memory/folders/{folder_id}/items')
def list_folder_items(folder_id: int):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('SELECT id, title, filename, substr(text_content,1,200) as preview, tags, created_at FROM mem_item WHERE folder_id=? ORDER BY id DESC', (folder_id,))
    rows = c.fetchall()
    conn.close()
    items = [{'id': r[0], 'title': r[1], 'filename': r[2], 'preview': r[3], 'tags': r[4], 'created_at': r[5]} for r in rows]
    return {'items': items}

@app.post('/api/memory/upload')
def memory_upload(folder_id: int = Form(...), file: UploadFile = File(...), tags: str = Form('')):
    # Save file and try to extract text for embedding
    upload_dir = os.path.join(os.path.dirname(__file__), 'memory_uploads')
    os.makedirs(upload_dir, exist_ok=True)
    filename = file.filename
    path = os.path.join(upload_dir, filename)
    with open(path, 'wb') as f:
        f.write(file.file.read())
    # Try to extract text based on extension (best-effort for PDFs)
    ext = os.path.splitext(filename)[1].lower()
    text = ''
    try:
        if ext == '.pdf':
            text = _extract_text_pdf_best_effort(path)
        elif ext in ('.docx', '.doc'):
            doc = docx.Document(path)
            text = '\n'.join([p.text for p in doc.paragraphs])
        elif ext in ('.xlsx', '.xls'):
            wb = openpyxl.load_workbook(path)
            for sh in wb.worksheets:
                for row in sh.iter_rows(values_only=True):
                    text += ' '.join([str(c) for c in row if c is not None]) + '\n'
        else:
            try:
                with open(path, 'r', encoding='utf-8') as fh:
                    text = fh.read()
            except Exception:
                text = ''
    except Exception as e:
        logging.warning(f"Failed to extract text from {filename}: {e}")
    # Prepare OCR diagnostics
    ocr_info = {
        'attempted': False,
        'engine': None,  # 'pdf2image+poppler' | 'pdfium-render' | 'image'
        'tesseract_cmd': getattr(pytesseract.pytesseract, 'tesseract_cmd', None) if pytesseract is not None else None,
        'tesseract_version': None,
        'poppler_path': POPPLER_PATH,
        'pdfium_available': pdfium is not None,
        'error': None,
    }
    try:
        if pytesseract is not None:
            ocr_info['tesseract_version'] = str(pytesseract.get_tesseract_version())
    except Exception:
        ocr_info['tesseract_version'] = None
    # If no text was extractable from parsers, try OCR for PDFs and images (robust)
    if (not text or not text.strip()) and pytesseract is not None:
        try:
            ocr_text = ''
            ocr_info['attempted'] = True
            if ext == '.pdf' and fitz is not None:
                # Prefer PyMuPDF rendering to avoid Poppler dependency
                ocr_text = _ocr_pdf_via_fitz(path, TESSERACT_LANGS)
                if ocr_text.strip():
                    ocr_info['engine'] = 'fitz-render'
                else:
                    ocr_info['error'] = (ocr_info.get('error') or '') or 'fitz render produced no text'
            if (not ocr_text or not ocr_text.strip()) and ext == '.pdf' and convert_from_path is not None:
                # Convert PDF pages to images and OCR each; allow POPPLER_PATH for Windows
                try:
                    kwargs = {'dpi': 200}
                    if POPPLER_PATH:
                        kwargs['poppler_path'] = POPPLER_PATH
                    pages = convert_from_path(path, **kwargs)
                    for img in pages:
                        ocr_text += pytesseract.image_to_string(img, lang=TESSERACT_LANGS) + '\n'
                    ocr_info['engine'] = 'pdf2image+poppler'
                except Exception as e:
                    logging.warning(f"[OCR] pdf2image conversion failed for {filename}: {e}")
                    ocr_info['error'] = f"pdf2image failed: {e}"
                    # Fallback to PDFium-based rendering if available
                    if pdfium is not None:
                        try:
                            pdf = pdfium.PdfDocument(path)
                            n_pages = len(pdf)
                            for i in range(n_pages):
                                page = pdf[i]
                                bitmap = page.render(scale=2.0, rotation=0)  # ~200 DPI
                                pil = bitmap.to_pil()
                                ocr_text += pytesseract.image_to_string(pil, lang=TESSERACT_LANGS) + '\n'
                            ocr_info['engine'] = 'pdfium-render'
                        except Exception as e2:
                            logging.warning(f"[OCR] PDFium fallback failed for {filename}: {e2}")
                            ocr_info['error'] = f"pdfium failed: {e2}"
            elif ext == '.pdf' and convert_from_path is None and pdfium is not None and (not ocr_text or not ocr_text.strip()):
                try:
                    pdf = pdfium.PdfDocument(path)
                    n_pages = len(pdf)
                    for i in range(n_pages):
                        page = pdf[i]
                        bitmap = page.render(scale=2.0, rotation=0)
                        pil = bitmap.to_pil()
                        ocr_text += pytesseract.image_to_string(pil, lang=TESSERACT_LANGS) + '\n'
                    ocr_info['engine'] = 'pdfium-render'
                except Exception as e:
                    logging.warning(f"[OCR] PDFium rendering failed for {filename}: {e}")
                    ocr_info['error'] = f"pdfium failed: {e}"
            elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp'):
                try:
                    img = Image.open(path)
                    ocr_text = pytesseract.image_to_string(img, lang=TESSERACT_LANGS)
                    ocr_info['engine'] = 'image'
                except Exception as e:
                    logging.warning(f"[OCR] PIL/Pytesseract failed for image {filename}: {e}")
                    ocr_info['error'] = f"image ocr failed: {e}"
            if ocr_text and ocr_text.strip():
                text = ocr_text
                logging.info(f"[OCR] Extracted text from {filename}, length={len(text)}")
        except Exception as e:
            logging.warning(f"[OCR] Fallback failed for {filename}: {e}")
            ocr_info['error'] = f"fallback failed: {e}"
    elif (not text or not text.strip()) and pytesseract is None:
        logging.warning(f"[OCR] pytesseract not available; skipping OCR for {filename}. Install Tesseract and pytesseract.")
        ocr_info['attempted'] = False
        ocr_info['error'] = 'pytesseract not available'
    # Guard: if the only 'text' looks like just the filename, treat as empty
    if text and text.strip():
        base_name = os.path.basename(filename)
        if text.strip().strip('\n\r\t ') == base_name:
            logging.info(f"[OCR] Detected text equals filename for {filename}; treating as empty.")
            text = ''
    # fallback title
    title = filename
    # Chunk the text for better retrieval and create embeddings per chunk
    chunks = []
    if text and text.strip():
        chunks = split_text(text, max_len=250)
    else:
        # Do NOT use filename as text content. Store a single empty-content record (no embedding).
        chunks = ['']

    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        for idx, chunk in enumerate(chunks):
            chunk_title = title
            if len(chunks) > 1 and chunk.strip():
                chunk_title = f"{title} (part {idx+1}/{len(chunks)})"
            emb_blob = None
            if chunk and chunk.strip():
                try:
                    emb_blob = embedder.encode(chunk).tobytes()
                except Exception as _e:
                    logging.warning(f"[Embed] Failed to encode chunk for {filename}: {_e}")
                    emb_blob = None
            c.execute('INSERT INTO mem_item (folder_id, title, filename, text_content, tags, embedding, created_at) VALUES (?, ?, ?, ?, ?, ?, strftime("%s", "now"))',
                      (folder_id, chunk_title, filename, chunk, tags, emb_blob))
        conn.commit()
    finally:
        conn.close()
    return {'status': 'ok', 'chunks': len(chunks), 'ocr': ocr_info}

@app.post('/api/memory/note')
async def memory_note(req: Request):
    # Expect JSON body: { folder_id: int, text: str, tags?: str }
    try:
        data = await req.json()
    except Exception:
        data = {}
    folder_id = data.get('folder_id') if isinstance(data, dict) else None
    text = data.get('text') if isinstance(data, dict) else None
    tags = data.get('tags') if isinstance(data, dict) else ''
    if not folder_id or not text:
        return JSONResponse({'error': 'folder_id and text required'}, status_code=400)
    # Chunk the note text into smaller pieces and store each with its own embedding
    chunks = split_text(text, max_len=250) if text else [text]
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        for idx, chunk in enumerate(chunks):
            title = None
            if len(chunks) > 1:
                title = f"note_part_{idx+1}"
            emb = embedder.encode(chunk)
            c.execute('INSERT INTO mem_item (folder_id, title, filename, text_content, tags, embedding, created_at) VALUES (?, ?, ?, ?, ?, ?, strftime("%s", "now"))',
                      (folder_id, title, None, chunk, tags, emb.tobytes()))
        conn.commit()
    finally:
        conn.close()
    return {'status': 'ok', 'chunks': len(chunks)}

@app.delete('/api/memory/item/{item_id}')
def delete_mem_item(item_id: int):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('DELETE FROM mem_item WHERE id=?', (item_id,))
    conn.commit()
    conn.close()
    return {'status': 'ok'}

@app.get('/api/memory/item/{item_id}/context')
def item_context(item_id: int):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('SELECT text_content, filename, tags FROM mem_item WHERE id=?', (item_id,))
    row = c.fetchone()
    conn.close()
    if not row:
        return JSONResponse({'error': 'not found'}, status_code=404)
    text, filename, tags = row
    return {'id': item_id, 'text': text, 'filename': filename, 'tags': tags}

@app.get('/api/memory/tag/{tag}/context')
def tag_context(tag: str):
    # Return concatenated small previews for items matching tag (efficient retrieval via DB + embeddings later)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    like = f'%{tag}%'
    c.execute('SELECT id, substr(text_content,1,1000), filename FROM mem_item WHERE tags LIKE ? OR filename LIKE ?', (like, like))
    rows = c.fetchall()
    conn.close()
    items = [{'id': r[0], 'preview': r[1], 'filename': r[2]} for r in rows]
    return {'items': items}


@app.get('/api/memory/suggest')
def memory_suggest(q: Optional[str] = None, limit: int = 10):
    """Return quick suggestions for folders, item tags, and filenames matching the prefix q.
    This endpoint is designed to be extremely fast (simple LIKE queries) and safe to call on every keystroke.
    """
    q = (q or '').strip()
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    suggestions = []
    try:
        if q:
            like = f'%{q}%'
            # folder names
            c.execute('SELECT DISTINCT name FROM mem_folder WHERE name LIKE ? LIMIT ?', (like, limit))
            folders = [r[0] for r in c.fetchall() if r[0]]
            # tags from mem_item.tags (comma separated) - simple LIKE match
            c.execute('SELECT DISTINCT tags FROM mem_item WHERE tags LIKE ? LIMIT ?', (like, limit))
            tags = []
            for row in c.fetchall():
                if not row[0]:
                    continue
                for t in str(row[0]).split(','):
                    tt = t.strip()
                    if q.lower() in tt.lower() and tt not in tags:
                        tags.append(tt)
                        if len(tags) >= limit:
                            break
                if len(tags) >= limit:
                    break
            # filenames
            c.execute('SELECT DISTINCT filename FROM mem_item WHERE filename LIKE ? LIMIT ?', (like, limit))
            files = [r[0] for r in c.fetchall() if r[0]]
        else:
            # when no query, return most recent folder names and tags
            c.execute('SELECT DISTINCT name FROM mem_folder ORDER BY id DESC LIMIT ?', (limit,))
            folders = [r[0] for r in c.fetchall() if r[0]]
            c.execute('SELECT DISTINCT tags FROM mem_item ORDER BY id DESC LIMIT ?', (limit,))
            tags = []
            for row in c.fetchall():
                if not row[0]:
                    continue
                for t in str(row[0]).split(','):
                    tt = t.strip()
                    if tt and tt not in tags:
                        tags.append(tt)
                        if len(tags) >= limit:
                            break
                if len(tags) >= limit:
                    break
            c.execute('SELECT DISTINCT filename FROM mem_item ORDER BY id DESC LIMIT ?', (limit,))
            files = [r[0] for r in c.fetchall() if r[0]]
        # Merge results keeping simple structure
        suggestions = {
            'folders': folders[:limit],
            'tags': tags[:limit],
            'files': files[:limit]
        }
    finally:
        conn.close()
    return suggestions


@app.post('/api/memory/search')
async def memory_search(request: Request):
    """Search mem_items by semantic similarity to a query. POST body: { query: str, top_k: int=5, tag: Optional[str] }
    Returns top_k items with preview and similarity score.
    """
    try:
        data = await request.json()
    except Exception:
        data = {}
    q = data.get('query') or ''
    top_k = int(data.get('top_k', 5))
    tag = data.get('tag')
    if not q:
        return JSONResponse({'error': 'query required'}, status_code=400)

    q_emb = embedder.encode(q)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    if tag:
        like = f'%{tag}%'
        c.execute('SELECT id, text_content, filename, tags, embedding FROM mem_item WHERE tags LIKE ? OR filename LIKE ?', (like, like))
    else:
        c.execute('SELECT id, text_content, filename, tags, embedding FROM mem_item')
    rows = c.fetchall()
    conn.close()

    import numpy as _np
    results = []
    for r in rows:
        try:
            emb_bytes = r[4]
            if not emb_bytes:
                continue
            emb = _np.frombuffer(emb_bytes, dtype=_np.float32)
            sim = float(_np.dot(q_emb, emb) / (_np.linalg.norm(q_emb) * _np.linalg.norm(emb)))
            results.append({'id': r[0], 'preview': (r[1] or '')[:1000], 'filename': r[2], 'tags': r[3], 'score': sim})
        except Exception:
            continue
    # sort by descending similarity
    results.sort(key=lambda x: x['score'], reverse=True)
    return {'items': results[:top_k]}


# Delete a memory folder and all its items
@app.delete('/api/memory/folders/{folder_id}')
def delete_mem_folder(folder_id: int):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    # delete items first
    c.execute('DELETE FROM mem_item WHERE folder_id=?', (folder_id,))
    c.execute('DELETE FROM mem_folder WHERE id=?', (folder_id,))
    conn.commit()
    conn.close()
    return {'status': 'ok'}


# Rename (or update tag for) a folder
class UpdateFolderRequest(BaseModel):
    name: Optional[str] = None
    tag: Optional[str] = None

@app.post('/api/memory/folders/{folder_id}/rename')
def rename_mem_folder(folder_id: int, req: UpdateFolderRequest):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    # Only update provided fields
    if req.name is not None and req.tag is not None:
        c.execute('UPDATE mem_folder SET name=?, tag=? WHERE id=?', (req.name, req.tag, folder_id))
    elif req.name is not None:
        c.execute('UPDATE mem_folder SET name=? WHERE id=?', (req.name, folder_id))
    elif req.tag is not None:
        c.execute('UPDATE mem_folder SET tag=? WHERE id=?', (req.tag, folder_id))
    conn.commit()
    conn.close()
    return {'status': 'ok'}


# Update an individual memory item (title / tags)
class UpdateItemRequest(BaseModel):
    title: Optional[str] = None
    tags: Optional[str] = None

@app.post('/api/memory/item/{item_id}/update')
def update_mem_item(item_id: int, req: UpdateItemRequest):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    if req.title is not None and req.tags is not None:
        c.execute('UPDATE mem_item SET title=?, tags=? WHERE id=?', (req.title, req.tags, item_id))
    elif req.title is not None:
        c.execute('UPDATE mem_item SET title=? WHERE id=?', (req.title, item_id))
    elif req.tags is not None:
        c.execute('UPDATE mem_item SET tags=? WHERE id=?', (req.tags, item_id))
    conn.commit()
    conn.close()
    return {'status': 'ok'}
# Endpoint to get chat history (list of chat_ids)
@app.get('/api/chats')
def get_chats(service: Optional[str] = None):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    if service:
        c.execute('SELECT DISTINCT chat_id, chat_name FROM chat WHERE service=? ORDER BY id DESC', (service,))
    else:
        c.execute('SELECT DISTINCT chat_id, chat_name, service FROM chat ORDER BY id DESC')
    rows = c.fetchall()
    conn.close()
    chats = []
    for r in rows:
        # rows can be (chat_id, chat_name) or (chat_id, chat_name, service)
        if len(r) == 2:
            chat_id, chat_name = r
            chats.append({"chat_id": chat_id, "chat_name": chat_name})
        else:
            chat_id, chat_name, svc = r
            chats.append({"chat_id": chat_id, "chat_name": chat_name, "service": svc})
    return {"chats": chats}

# Endpoint to get messages for a chat_id
@app.get('/api/chat/{chat_id}')
def get_chat_messages(chat_id: str, service: Optional[str] = None):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    if service:
        c.execute('SELECT id, user, llm, timestamp, doc_info, tags FROM chat WHERE chat_id=? AND service=? ORDER BY id ASC', (chat_id, service))
    else:
        c.execute('SELECT id, user, llm, timestamp, doc_info, chat_name, tags FROM chat WHERE chat_id=? ORDER BY id ASC', (chat_id,))
    messages = []
    chat_name = None
    for row in c.fetchall():
        # If chat_name is included in the row tuple, pick it up; else leave None
        if len(row) == 6:
            row_id, user_text, llm_text, timestamp, doc_info, tags = row
        else:
            row_id, user_text, llm_text, timestamp, doc_info, chat_name, tags = row
        if user_text:
            messages.append({"sender": "user", "text": user_text, "timestamp": timestamp, "doc_info": doc_info, "row_id": row_id, "tags": tags})
        if llm_text:
            messages.append({"sender": "ai", "text": llm_text, "timestamp": timestamp, "doc_info": doc_info, "row_id": row_id, "tags": tags})
    conn.close()
    return {"messages": messages, "chat_name": chat_name}

# -------------------- Chat management: delete / rename --------------------
class DeleteChatRequest(BaseModel):
    service: Optional[str] = None

@app.delete('/api/chat/{chat_id}')
def delete_chat(chat_id: str, service: Optional[str] = None):
    """Delete all chat rows for a given chat_id (optionally scoped to service). Also clears chat_state.

    File deletion policy: ONLY delete the file pointed to by chat_state.doc_path for this chat/service.
    No other files are removed.
    """
    conn = sqlite3.connect(DB_PATH)
    try:
        c = conn.cursor()
        # Look up doc_path from chat_state BEFORE deleting rows
        state = _get_chat_state(conn, chat_id, service)
        doc_path = state.get('doc_path') if isinstance(state, dict) else None

        # Delete DB rows
        if service:
            c.execute('DELETE FROM chat WHERE chat_id=? AND service=?', (chat_id, service))
            c.execute('DELETE FROM chat_state WHERE chat_id=? AND service=?', (chat_id, service))
            c.execute('DELETE FROM chat_artifact WHERE chat_id=? AND service=?', (chat_id, service))
        else:
            c.execute('DELETE FROM chat WHERE chat_id=?', (chat_id,))
            c.execute('DELETE FROM chat_state WHERE chat_id=?', (chat_id,))
            c.execute('DELETE FROM chat_artifact WHERE chat_id=?', (chat_id,))
        conn.commit()

        # Delete ONLY the doc_path file (safe-guard to managed folders)
        deleted = []
        failed = []
        if doc_path:
            import os as _os
            from pathlib import Path as _P
            p = _P(doc_path)
            # restrict to our managed folders
            try:
                home = _P.home()
                docs = home / 'Documents'
                allowed1 = (docs / 'SarvajnaGPT').resolve()
            except Exception:
                allowed1 = None
            try:
                repo_root = _P(__file__).resolve().parent.parent
                allowed2 = (repo_root / 'SarvajnaGPT').resolve()
            except Exception:
                allowed2 = None
            def _is_allowed(path: _P) -> bool:
                try:
                    rp = path.resolve()
                    ok1 = bool(allowed1 and str(rp).startswith(str(allowed1)))
                    ok2 = bool(allowed2 and str(rp).startswith(str(allowed2)))
                    return ok1 or ok2
                except Exception:
                    return False
            try:
                if p.is_file() and _is_allowed(p):
                    _os.remove(str(p))
                    deleted.append(str(p))
            except Exception as e:
                failed.append({"path": str(p), "error": str(e)})

        return {"status": "ok", "deleted_chat_id": chat_id, "service": service, **({"deleted_files": deleted} if deleted else {}), **({"failed_files": failed} if failed else {})}
    finally:
        conn.close()

# -------------------- Chat State (Power Mode persistence) --------------------
class ChatStateUpsert(BaseModel):
    chat_id: str
    service: Optional[str] = None
    persistent_tags: Optional[List[str]] = None
    doc_path: Optional[str] = None

def _get_chat_state(conn: sqlite3.Connection, chat_id: str, service: Optional[str] = None):
    c = conn.cursor()
    if service:
        c.execute('SELECT id, persistent_tags, doc_path FROM chat_state WHERE chat_id=? AND service=? ORDER BY id DESC LIMIT 1', (chat_id, service))
    else:
        c.execute('SELECT id, persistent_tags, doc_path FROM chat_state WHERE chat_id=? ORDER BY id DESC LIMIT 1', (chat_id,))
    row = c.fetchone()
    if not row:
        return None
    cid, ptags, dpath = row
    tags = []
    if ptags:
        try:
            tags = [t.strip() for t in str(ptags).split(',') if t and str(t).strip()]
        except Exception:
            tags = []
    return {"id": cid, "persistent_tags": tags, "doc_path": dpath}

def _upsert_chat_state(conn: sqlite3.Connection, chat_id: str, service: Optional[str], tags: Optional[List[str]], doc_path: Optional[str]):
    c = conn.cursor()
    existing = _get_chat_state(conn, chat_id, service)
    tags_str = None
    if tags:
        tags_norm = [str(t).strip() for t in tags if t and str(t).strip()]
        if tags_norm:
            tags_str = ','.join(tags_norm)
    now = int(time.time())
    if existing:
        # merge behavior: overwrite non-null inputs; keep previous when inputs are None
        new_tags_str = tags_str if tags_str is not None else ','.join(existing.get('persistent_tags') or []) or None
        new_doc = doc_path if doc_path is not None else existing.get('doc_path')
        c.execute('UPDATE chat_state SET persistent_tags=?, doc_path=?, updated_at=? WHERE id=?', (new_tags_str, new_doc, now, existing['id']))
    else:
        c.execute('INSERT INTO chat_state (chat_id, service, persistent_tags, doc_path, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?)', (chat_id, service, tags_str, doc_path, now, now))
    conn.commit()

@app.post('/api/chat_state')
def set_chat_state(req: ChatStateUpsert):
    conn = sqlite3.connect(DB_PATH)
    try:
        _upsert_chat_state(conn, req.chat_id, req.service, req.persistent_tags, req.doc_path)
        state = _get_chat_state(conn, req.chat_id, req.service)
        return {"status": "ok", "state": state}
    finally:
        conn.close()

@app.get('/api/chat_state/{chat_id}')
def get_chat_state(chat_id: str, service: Optional[str] = None):
    conn = sqlite3.connect(DB_PATH)
    try:
        state = _get_chat_state(conn, chat_id, service)
        return {"state": state}
    finally:
        conn.close()

# -------------------- Global: Load TTS once --------------------
device = "cuda" if torch.cuda.is_available() else "cpu"
MODEL_NAME = "tts_models/multilingual/multi-dataset/xtts_v2"
tts_engine = TTS(model_name=MODEL_NAME)
tts_engine.to(device)
# Print available speakers and languages for debugging
try:
    speakers = tts_engine.speakers if hasattr(tts_engine, 'speakers') else []
    languages = tts_engine.languages if hasattr(tts_engine, 'languages') else []
    logging.info(f"[TTS] Available speakers: {speakers}")
    logging.info(f"[TTS] Available languages: {languages}")
except Exception as e:
    logging.warning(f"[TTS] Could not list speakers/languages: {e}")

LANG_SPEAKER_MAP = {
    'en': 'Andrew Chipper',
    'hi': 'Kumar Dahl',
    'fr': 'Gracie Wise',
    'de': 'Ludvig Milivoj',
    'es': 'Marcos Rudaski',
    'it': 'Eugenio Mataracı',
    'ru': 'Lidiya Szekeres',
    'zh-cn': 'Xavier Hayasaka',
    'ja': 'Kazuhiko Atallah',
    'ko': 'Nova Hogarth',
    'pt': 'Luis Moray',
    'ar': 'Suad Qasim',
    'tr': 'Asya Anara',
    'uk': 'Szofi Granger',
    'pl': 'Zofija Kendrick',
    'nl': 'Camilla Holmström',
    'bn': 'Chandra MacFarland',
    'ta': 'Narelle Moon',
    'te': 'Alma María',
}
DEFAULT_SPEAKER = LANG_SPEAKER_MAP.get('en', 'Kumar Dahl')
TTS_SAMPLE_RATE = tts_engine.synthesizer.output_sample_rate if hasattr(tts_engine, 'synthesizer') else 24000
logging.info(f"[TTS] Model loaded once at startup ({MODEL_NAME}) on {device}")
if not DEFAULT_SPEAKER:
    logging.warning("[TTS] No speakers found for the model. TTS might fail or use a default voice.")

# -------------------- Global: Language Detector --------------------
supported_langs = [
    Language.ENGLISH, Language.HINDI, Language.FRENCH, Language.GERMAN,
    Language.SPANISH, Language.ITALIAN, Language.RUSSIAN, Language.CHINESE,
    Language.JAPANESE, Language.KOREAN, Language.PORTUGUESE, Language.ARABIC,
    Language.TURKISH, Language.UKRAINIAN, Language.POLISH, Language.DUTCH,
    Language.BENGALI, Language.TAMIL, Language.TELUGU
]
detector = LanguageDetectorBuilder.from_languages(*supported_langs).build()

# -------------------- Helpers --------------------
def clean_text(text: str) -> str:
    """Normalize unicode, keep Hindi + normal punctuation, remove emojis/markdown/symbols, and clean whitespace."""
    # Normalize text
    text = unicodedata.normalize("NFC", text)

    # Remove markdown-style asterisks and quotes
    text = re.sub(r'[*"]', '', text)

    # Remove numbers for non-English (to avoid TTS num2words crash)
    # (Assume English if only ASCII, else remove numbers)
    if not re.match(r'^[\u0000-\u007F]*$', text):
        text = re.sub(r'\d+', '', text)

    # Remove emojis and other symbols outside common ranges (keep Devanagari, Latin, punctuation)
    text = re.sub(r'[^\w\s\u0900-\u097F.,!?()\-:]', '', text, flags=re.UNICODE)

    # Collapse multiple spaces/newlines
    text = re.sub(r'\s*\n\s*', '\n', text)   # normalize newlines
    text = re.sub(r'[ \t]+', ' ', text)      # collapse spaces/tabs
    text = text.strip()

    return text

def split_text(text: str, max_len=150):
    words = text.split()
    chunks, current = [], []
    total_len = 0
    for w in words:
        if total_len + len(w) + 1 > max_len:
            chunks.append(" ".join(current))
            current = [w]
            total_len = len(w)
        else:
            current.append(w)
            total_len += len(w) + 1
    if current:
        chunks.append(" ".join(current))
    return chunks

async def synthesize_chunk_async(text: str, lang_code: str):
    """Generate speech for one chunk in a thread-safe manner."""
    loop = asyncio.get_running_loop()
    
    import traceback

    def tts_sync():
        # Pick speaker based on language
        speaker = LANG_SPEAKER_MAP.get(lang_code, DEFAULT_SPEAKER)
        kwargs = {'text': text, 'language': lang_code, 'speaker': speaker}
        logging.info(f"[TTS] Calling tts_engine.tts with: {kwargs}")
        try:
            return tts_engine.tts(**kwargs)
        except Exception as e:
            logging.error(f"[TTS] Error in tts_engine.tts: {e}\n" + traceback.format_exc())
            raise

    # Offload the blocking TTS synthesis to a thread pool
    audio_array = await loop.run_in_executor(None, tts_sync)
    return np.array(audio_array)

def save_and_schedule_delete(wav_array, filename, sample_rate, delay=60):
    """Save wav and schedule auto deletion."""
    temp_dir = tempfile.gettempdir()
    wav_path = os.path.join(temp_dir, filename)
    static_dir = os.path.join(os.path.dirname(__file__), "static")
    os.makedirs(static_dir, exist_ok=True)
    out_path = os.path.join(static_dir, filename)

    sf.write(wav_path, wav_array, sample_rate)
    shutil.copyfile(wav_path, out_path)

    def _delete():
        time.sleep(delay)
        try:
            os.remove(out_path)
            os.remove(wav_path) # Also remove from temp
        except Exception:
            pass
    threading.Thread(target=_delete, daemon=True).start()

    return f"/static/{filename}"


# -------------------- REST: Chunked TTS Endpoint for Polling --------------------
@app.post("/api/voice_chunk")
async def voice_chunk_endpoint(request: Request):
    data = await request.json()
    text = data.get("text", "")
    chunk_index = int(data.get("chunk", 0))
    session_id = data.get("session_id") or str(uuid.uuid4())
    lang_code = data.get("lang", None)
    supported_langs = ['en', 'es', 'fr', 'de', 'it', 'pt', 'pl', 'tr', 'ru', 'nl', 'cs', 'ar', 'zh-cn', 'hu', 'ko', 'ja', 'hi']
    if not lang_code or lang_code == 'auto':
        try:
            detected_lang = detector.detect_language_of(text)
            # Map lingua detection to supported TTS code
            lingua_to_tts = {
                'en': 'en', 'hi': 'hi', 'fr': 'fr', 'de': 'de', 'es': 'es', 'it': 'it', 'ru': 'ru', 'zh': 'zh-cn',
                'ko': 'ko', 'ja': 'ja', 'ar': 'ar', 'pt': 'pt', 'pl': 'pl', 'nl': 'nl', 'cs': 'cs', 'tr': 'tr', 'hu': 'hu'
            }
            lingua_code = detected_lang.iso_code_639_1.name.lower() if detected_lang else 'en'
            lang_code = lingua_to_tts.get(lingua_code, 'en')
        except Exception:
            lang_code = 'en'
    # Fallback if unsupported
    if lang_code not in supported_langs:
        lang_code = 'en'
    processed_text = clean_text(text)
    chunks = split_text(processed_text, max_len=100)
    if chunk_index >= len(chunks):
        return JSONResponse({"error": "Chunk index out of range", "done": True, "chunks": []})
    chunk = chunks[chunk_index]
    import traceback
    try:
        audio_array = await synthesize_chunk_async(chunk, lang_code)
        duration_sec = len(audio_array) / TTS_SAMPLE_RATE
        filename = f"tts_{session_id}_{chunk_index}_{uuid.uuid4().hex}.wav"
        url = save_and_schedule_delete(audio_array, filename, TTS_SAMPLE_RATE, delay=600)
        # Always return a 'chunks' array for frontend compatibility
        return JSONResponse({
            "voiceUrl": url,
            "chunk": chunk_index,
            "durationSec": duration_sec,
            "lang": lang_code,
            "voice": DEFAULT_SPEAKER,
            "done": chunk_index == len(chunks) - 1,
            "chunks": [url],
            "debug": f"Processed: {chunk[:50]}..."
        })
    except Exception as e:
        logging.error(f"[TTS] Exception in voice_chunk_endpoint: {e}\n" + traceback.format_exc())
        return JSONResponse({"error": str(e), "chunk": chunk_index, "done": False, "traceback": traceback.format_exc(), "chunks": []}, status_code=500)

# static mount
app.mount("/static", StaticFiles(directory=os.path.join(os.path.dirname(__file__), "static")), name="static")

# Serve built frontend (if present) so LAN users only need backend port
_DIST_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'frontend', 'dist'))
if os.path.isdir(_DIST_DIR):
    try:
        app.mount("/", StaticFiles(directory=_DIST_DIR, html=True), name="frontend")
        logging.info(f"[Frontend] Serving built frontend from {_DIST_DIR}")
    except Exception as fe:
        logging.warning(f"[Frontend] Failed to mount dist: {fe}")


@app.post("/generate")
async def generate(request: Request):
    data = await request.json()
    prompt = data.get("prompt", "")
    try:
        response = generate_response(prompt)
        return JSONResponse({"response": response})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/summarize")
async def summarize(request: Request):
    data = await request.json()
    text = data.get("text", "")
    try:
        summary = summarize_text(text)
        return JSONResponse({"summary": summary})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/answer")
async def answer(request: Request):
    data = await request.json()
    question = data.get("question", "")
    context = data.get("context", "")
    try:
        answer = answer_question(question, context)
        return JSONResponse({"answer": answer})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/analyze_pdf")
async def analyze_pdf_endpoint(file: UploadFile = File(...)):
    try:
        result = analyze_pdf(file)
        return JSONResponse({"result": result})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/settings")
async def get_app_settings():
    settings = get_settings()
    return JSONResponse(settings)

@app.post("/settings")
async def update_app_settings(request: Request):
    data = await request.json()
    set_settings(data)
    return JSONResponse({"status": "updated"})

if __name__ == "__main__":
    import uvicorn, argparse
    parser = argparse.ArgumentParser(description="SarvajnaGPT Backend Server")
    parser.add_argument("--host", default="0.0.0.0")
    parser.add_argument("--port", type=int, default=8000)
    parser.add_argument("--reload", action="store_true")
    args = parser.parse_args()
    logging.info(f"Starting server on {args.host}:{args.port} (reload={args.reload})")
    uvicorn.run("main:app", host=args.host, port=args.port, reload=args.reload)

# Extract text from uploaded document
@app.post('/api/extract_text')
async def extract_text(file: UploadFile = File(...)):
    ext = file.filename.split('.')[-1].lower()
    content = ''
    try:
        if ext == 'pdf':
            # write to temp for consistent handling
            await file.seek(0)
            data = await file.read()
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            try:
                content = _extract_text_pdf_best_effort(tmp_path)
                if (not content or not content.strip()) and pytesseract is not None:
                    # Prefer fitz render
                    content = _ocr_pdf_via_fitz(tmp_path, TESSERACT_LANGS)
                    # pdf2image + Poppler
                    if (not content or not content.strip()) and convert_from_path is not None:
                        try:
                            kwargs = {'dpi': 200}
                            if POPPLER_PATH:
                                kwargs['poppler_path'] = POPPLER_PATH
                            images = convert_from_path(tmp_path, **kwargs)
                            ocr_text = ''
                            for im in images:
                                ocr_text += pytesseract.image_to_string(im, lang=TESSERACT_LANGS) + '\n'
                            content = ocr_text
                        except Exception:
                            pass
                    # PDFium
                    if (not content or not content.strip()) and pdfium is not None:
                        try:
                            pdf = pdfium.PdfDocument(tmp_path)
                            n_pages = len(pdf)
                            ocr_text = ''
                            for i in range(n_pages):
                                page = pdf[i]
                                bitmap = page.render(scale=2.0, rotation=0)
                                pil = bitmap.to_pil()
                                ocr_text += pytesseract.image_to_string(pil, lang=TESSERACT_LANGS) + '\n'
                            content = ocr_text
                        except Exception:
                            pass
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        elif ext == 'docx' and docx:
            doc = docx.Document(io.BytesIO(await file.read()))
            content = '\n'.join([p.text for p in doc.paragraphs])
        elif ext in ['xlsx', 'xls'] and openpyxl:
            wb = openpyxl.load_workbook(io.BytesIO(await file.read()), read_only=True)
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    content += '\t'.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
        elif ext == 'pptx' and pptx:
            pres = pptx.Presentation(io.BytesIO(await file.read()))
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'
        else:
            # fallback: try to read as text
            content = (await file.read()).decode(errors='ignore')
        # OCR fallback when nothing readable was found (images)
        if (not content or not content.strip()) and pytesseract is not None:
            try:
                await file.seek(0)
                data = await file.read()
                if ext in ['png','jpg','jpeg','bmp','tif','tiff']:
                    img = Image.open(io.BytesIO(data))
                    content = pytesseract.image_to_string(img, lang=TESSERACT_LANGS)
            except Exception as _e:
                logging.warning(f"[OCR] extract_text fallback failed for {file.filename}: {_e}")
        elif (not content or not content.strip()) and pytesseract is None:
            logging.warning(f"[OCR] pytesseract not available; skipping OCR for {file.filename}. Install Tesseract and pytesseract.")
    except Exception as e:
        content = f'[Error extracting text: {str(e)}]'
    # Guard: if OCR result equals the filename, treat as empty
    try:
        base_name = os.path.basename(file.filename)
        if content and content.strip() == base_name:
            content = ''
    except Exception:
        pass
    return {"text": content}

# --------- OCR Health Check ---------
@app.get('/api/ocr/health')
def ocr_health():
    """Report availability of OCR-related components to help diagnose issues."""
    def _tess_version():
        try:
            if pytesseract is None:
                return None
            return str(pytesseract.get_tesseract_version())
        except Exception:
            return None
    def _engine_present():
        try:
            cmd = getattr(pytesseract.pytesseract, 'tesseract_cmd', None) if pytesseract is not None else None
            if not cmd:
                return False
            # If it's a bare name like 'tesseract' on Windows, accept it only if we can get a version
            if os.name == 'nt' and not os.path.isabs(cmd):
                try:
                    v = str(pytesseract.get_tesseract_version()) if pytesseract is not None else None
                    return v is not None and len(v) > 0
                except Exception:
                    return False
            return os.path.isfile(cmd) if os.path.isabs(cmd) else True
        except Exception:
            return False
    return {
        'pytesseract_available': pytesseract is not None,
        'tesseract_cmd': getattr(pytesseract.pytesseract, 'tesseract_cmd', None) if pytesseract is not None else None,
        'tesseract_version': _tess_version(),
        'tesseract_langs': TESSERACT_LANGS,
        'pdf2image_available': convert_from_path is not None,
        'poppler_path': POPPLER_PATH,
        'poppler_optional': True,  # Explicit: Poppler is optional; PyMuPDF/PDFium cover rendering
        'pdfium_available': pdfium is not None,
        'pillow_available': True,
        'fitz_available': fitz is not None,
        'pdfplumber_available': pdfplumber is not None,
        'notes': 'Poppler is optional. Preferred stack order: text extraction (PyMuPDF/pdfplumber/pdfminer/PyPDF2) then OCR via PyMuPDF render -> pdf2image(poppler) -> PDFium. If PyMuPDF works, pdf2image/poppler is skipped.',
        'engine_present': _engine_present(),
    }
    
@app.post('/api/ocr/refresh')
def ocr_refresh():
    """Re-run Tesseract/Poppler auto-detection to heal environment after installs (Windows)."""
    changed = _ensure_tesseract_and_poppler()
    version = None
    try:
        if pytesseract is not None:
            version = str(pytesseract.get_tesseract_version())
    except Exception:
        version = None
    return {
        'changed': changed,
        'tesseract_cmd': getattr(pytesseract.pytesseract, 'tesseract_cmd', None) if pytesseract is not None else None,
        'tesseract_version': version,
        'poppler_path': POPPLER_PATH,
    }